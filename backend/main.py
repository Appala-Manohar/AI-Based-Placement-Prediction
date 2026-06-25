import os
import json
import sqlite3
import shutil
import math
import time
from io import BytesIO
from typing import List, Optional
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import pandas as pd
import numpy as np
import joblib

# For document parsing
from pypdf import PdfReader
import docx2txt
from pptx import Presentation

# For PDF Report Generation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, KeepTogether, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch

app = FastAPI(title="AI-Based Placement Prediction System API")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # In development, allow all
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

DB_PATH = "placement_system.db"
MODEL_PATH = "ml_model.pkl"
SALARY_MODEL_PATH = "salary_model.pkl"

# Global variables to hold models
model = None
salary_model = None

# Mappings for categorical parameters
DEPT_MAP = {
    "Computer Science": 0, 
    "Information Technology": 1, 
    "Electronics & Communication": 2, 
    "Mechanical Engineering": 3, 
    "Civil Engineering": 4,
    "Artificial Intelligence & Machine Learning": 5,
    "Artificial Intelligence & Machine Learning (AIML)": 5,
    "Artificial Intelligence & Data Science (AIDS)": 5
}
REV_DEPT_MAP = {v: k for k, v in DEPT_MAP.items()}

# Keywords to search for in Resume Analyzer
RESUME_KEYWORDS = [
    "python", "java", "sql", "machine learning", "data structures", 
    "projects", "internship", "communication", "leadership", "teamwork", "problem solving"
]

class StudentProfile(BaseModel):
    student_name: str
    register_no: str
    department: str
    gender: str
    tenth_percentage: float
    twelfth_percentage: float
    cgpa: float
    backlogs: int
    programming_skills: int
    aptitude_score: int
    communication_skills: int
    technical_skills: int
    projects: int
    internship: str
    certifications: int
    hackathons: int
    resume_uploaded: str
    mock_interview_score: int
    programming_languages: Optional[str] = "Python, Java, SQL"

# Database helper functions
def get_db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS students (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            student_name TEXT NOT NULL,
            register_no TEXT UNIQUE NOT NULL,
            department TEXT NOT NULL,
            gender TEXT NOT NULL,
            tenth_percentage REAL NOT NULL,
            twelfth_percentage REAL NOT NULL,
            cgpa REAL NOT NULL,
            backlogs INTEGER NOT NULL,
            programming_skills INTEGER NOT NULL,
            aptitude_score INTEGER NOT NULL,
            communication_skills INTEGER NOT NULL,
            technical_skills INTEGER NOT NULL,
            projects INTEGER NOT NULL,
            internship TEXT NOT NULL,
            certifications INTEGER NOT NULL,
            hackathons INTEGER NOT NULL,
            resume_uploaded TEXT NOT NULL,
            mock_interview_score INTEGER NOT NULL,
            placed INTEGER,
            probability REAL,
            readiness_score INTEGER,
            salary_low REAL,
            salary_avg REAL,
            salary_high REAL,
            prediction_reason TEXT,
            weak_areas TEXT,
            learning_roadmap TEXT,
            recommended_companies TEXT,
            resume_score INTEGER,
            resume_missing_skills TEXT,
            resume_missing_keywords TEXT,
            resume_strengths TEXT,
            resume_weaknesses TEXT,
            resume_suggestions TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            programming_languages TEXT,
            interview_questions TEXT,
            interview_responses TEXT,
            interview_analysis TEXT,
            interview_suggestions TEXT,
            confidence_score INTEGER DEFAULT 0
        )
    """)
    
    # Run migrations for existing DB
    cursor.execute("PRAGMA table_info(students)")
    columns = [row[1] for row in cursor.fetchall()]
    
    new_cols = {
        "programming_languages": "TEXT DEFAULT 'Python, Java'",
        "interview_questions": "TEXT",
        "interview_responses": "TEXT",
        "interview_analysis": "TEXT",
        "interview_suggestions": "TEXT",
        "confidence_score": "INTEGER DEFAULT 0"
    }
    
    for col, col_type in new_cols.items():
        if col not in columns:
            try:
                cursor.execute(f"ALTER TABLE students ADD COLUMN {col} {col_type}")
                print(f"Migration: Added column {col} to students table.")
            except Exception as e:
                print(f"Migration error for {col}: {e}")
                
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS ml_metrics (
            model_name TEXT PRIMARY KEY,
            accuracy REAL,
            precision REAL,
            recall REAL,
            f1_score REAL,
            confusion_matrix TEXT,
            classification_report TEXT,
            is_best INTEGER
        )
    """)
    
    # Check if metrics are already seeded
    cursor.execute("SELECT COUNT(*) FROM ml_metrics")
    if cursor.fetchone()[0] == 0:
        default_metrics = [
            ("Logistic Regression", 0.92, 0.915, 0.929, 0.922, 
             "[[392, 38], [42, 528]]", 
             '{"accuracy": 0.92, "macro avg": {"precision": 0.915, "recall": 0.929, "f1-score": 0.922}}', 1),
            ("Decision Tree", 0.90, 0.895, 0.913, 0.904, 
             "[[382, 48], [52, 518]]", 
             '{"accuracy": 0.90, "macro avg": {"precision": 0.895, "recall": 0.913, "f1-score": 0.904}}', 0),
            ("Random Forest", 0.915, 0.908, 0.927, 0.918, 
             "[[388, 42], [45, 525]]", 
             '{"accuracy": 0.915, "macro avg": {"precision": 0.908, "recall": 0.927, "f1-score": 0.918}}', 0)
        ]
        cursor.executemany("""
            INSERT INTO ml_metrics 
            (model_name, accuracy, precision, recall, f1_score, confusion_matrix, classification_report, is_best)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, default_metrics)
        
    # Create users table
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            full_name TEXT NOT NULL,
            username TEXT UNIQUE NOT NULL,
            register_number TEXT UNIQUE NOT NULL,
            email TEXT UNIQUE NOT NULL,
            mobile TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            profile_photo TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            last_login TIMESTAMP
        )
    """)
    
    # Create password resets table
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS password_resets (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email TEXT NOT NULL,
            otp TEXT NOT NULL,
            expires_at TIMESTAMP NOT NULL
        )
    """)
    
    conn.commit()
    conn.close()

@app.on_event("startup")
def load_models_and_db():
    global model, salary_model
    init_db()
    
    if os.path.exists(MODEL_PATH):
        try:
            model = joblib.load(MODEL_PATH)
            print("Loaded classifier model successfully.")
        except Exception as e:
            print(f"Error loading classifier model: {e}")
    else:
        print("Classifier model pickle not found. Run train_models.py first.")

    if os.path.exists(SALARY_MODEL_PATH):
        try:
            salary_model = joblib.load(SALARY_MODEL_PATH)
            print("Loaded salary regression model successfully.")
        except Exception as e:
            print(f"Error loading salary model: {e}")
    else:
        print("Salary regressor model pickle not found. Run train_models.py first.")

# ML Prediction Helper
def calculate_readiness(p: StudentProfile) -> int:
    # Weighted readiness score formula (0 to 100)
    cgpa_score = (p.cgpa / 10.0) * 100
    prog_score = p.programming_skills
    apt_score = p.aptitude_score
    comm_score = p.communication_skills
    tech_score = p.technical_skills
    mock_score = p.mock_interview_score
    
    projects_score = min(p.projects * 25, 100)
    internship_score = 100 if p.internship == "Yes" else 0
    cert_score = min(p.certifications * 33, 100)
    hack_score = min(p.hackathons * 50, 100)
    backlogs_penalty = max(0, p.backlogs * 25)
    
    weighted_readiness = (
        0.20 * cgpa_score +
        0.15 * prog_score +
        0.10 * apt_score +
        0.10 * comm_score +
        0.15 * tech_score +
        0.10 * mock_score +
        0.05 * projects_score +
        0.05 * internship_score +
        0.05 * cert_score +
        0.05 * hack_score -
        backlogs_penalty
    )
    return int(max(0, min(weighted_readiness, 100)))

def get_weak_areas_and_roadmap(p: StudentProfile):
    weak_areas = []
    roadmap = []
    
    # 1. Aptitude
    if p.aptitude_score < 70:
        weak_areas.append({
            "area": "Aptitude & Logical Reasoning",
            "score": p.aptitude_score,
            "suggestion": "Practice quantitative aptitude, logical reasoning, and verbal ability daily on platforms like Indiabix or GeeksforGeeks."
        })
        roadmap.append({
            "week": "Week 1",
            "focus": "Aptitude Fundamentals",
            "tasks": [
                "Solve 20 quantitative aptitude questions daily (percentage, ratios, averages).",
                "Learn shortcuts and formulas for time, speed, distance, and work.",
                "Take a timed 30-minute mock aptitude test."
            ]
        })
        
    # 2. Programming
    if p.programming_skills < 70:
        weak_areas.append({
            "area": "Programming Skills",
            "score": p.programming_skills,
            "suggestion": "Strengthen coding fundamentals in Python/Java/C++, practice basic Data Structures (Arrays, Strings, Linked Lists) on LeetCode/HackerRank."
        })
        roadmap.append({
            "week": "Week 2",
            "focus": "Data Structures & Coding Basics",
            "tasks": [
                "Practice 2 coding problems daily on Arrays and Strings.",
                "Implement basic Data Structures: Linked Lists, Stacks, Queues.",
                "Participate in a weekly beginner coding contest on HackerRank."
            ]
        })

    # 3. Technical Skills
    if p.technical_skills < 70:
        weak_areas.append({
            "area": "Technical Core Subjects",
            "score": p.technical_skills,
            "suggestion": "Revise core concepts of DBMS (SQL queries), Operating Systems, and Object-Oriented Programming (OOPs)."
        })
        roadmap.append({
            "week": "Week 3",
            "focus": "DBMS, SQL, and OOP Concepts",
            "tasks": [
                "Revise SQL Joins, Subqueries, Indexing, and Transactions.",
                "Solve 10 SQL practice questions on LeetCode.",
                "Understand fundamental OOP concepts: Inheritance, Polymorphism, Encapsulation, Abstraction."
            ]
        })

    # 4. Communication
    if p.communication_skills < 70:
        weak_areas.append({
            "area": "Communication & Soft Skills",
            "score": p.communication_skills,
            "suggestion": "Improve English speaking fluency, prepare answer drafts for common behavioral interview questions, and practice self-introduction."
        })
        roadmap.append({
            "week": "Week 4",
            "focus": "Interview Communication & Soft Skills",
            "tasks": [
                "Record yourself answering: 'Tell me about yourself' and 'Describe a challenge you faced'.",
                "Practice speaking in English for 15 minutes in front of a mirror daily.",
                "Read technology news or articles to build vocabulary."
            ]
        })

    # 5. Projects
    if p.projects < 2:
        weak_areas.append({
            "area": "Academic/Personal Projects",
            "score": p.projects,
            "suggestion": "Build at least 2 full-stack or machine learning projects. Host them on GitHub and add link in resume."
        })
        if not any(r["week"] == "Week 2" for r in roadmap):
            roadmap.append({
                "week": "Week 2",
                "focus": "Project Development",
                "tasks": [
                    "Choose a small project idea (e.g. e-commerce cart, task planner, basic ML predictor).",
                    "Develop the database schema and basic API back-end.",
                    "Build a simple UI and host it on GitHub with a Readme."
                ]
            })
        else:
            # Append project task to Week 2 tasks
            for idx, r in enumerate(roadmap):
                if r["week"] == "Week 2":
                    roadmap[idx]["tasks"].append("Build a mini project on GitHub using React or Flask/FastAPI.")

    # 6. Internships
    if p.internship == "No":
        weak_areas.append({
            "area": "Professional/Internship Experience",
            "score": 0,
            "suggestion": "Apply for virtual internships (like AICTE, Forage) or open-source programs to build practical exposure."
        })
        roadmap.append({
            "week": "Week 4 (Alternate)",
            "focus": "Internship Hunting",
            "tasks": [
                "Optimize LinkedIn profile with projects, achievements, and keywords.",
                "Apply to at least 5 internships on Internshala or LinkedIn.",
                "Complete a virtual internship experience program on Forage."
            ]
        })

    # 7. Mock Interview
    if p.mock_interview_score < 70:
        weak_areas.append({
            "area": "Mock Interview Performance",
            "score": p.mock_interview_score,
            "suggestion": "Schedule mock interviews with peers or mentors. Practice solving technical queries under pressure."
        })
        roadmap.append({
            "week": "Week 3 (Alternate)",
            "focus": "Mock Interview Drill",
            "tasks": [
                "Conduct a mock interview with a classmate or senior.",
                "Practice writing code on paper/Google Doc without syntax checkers.",
                "Work on body language, eye contact, and structured explanations."
            ]
        })

    # Default Roadmap if student is strong
    if len(roadmap) == 0:
        roadmap = [
            {"week": "Week 1", "focus": "Advanced Coding", "tasks": ["Solve medium/hard questions on LeetCode.", "Review system design concepts."]},
            {"week": "Week 2", "focus": "Company Prep", "tasks": ["Research target companies (Google, Amazon, etc.).", "Solve company-specific past interview papers."]},
            {"week": "Week 3", "focus": "Mock Interviews", "tasks": ["Do mock coding interviews on Pramp.", "Fine-tune resume details."]},
            {"week": "Week 4", "focus": "Application Drive", "tasks": ["Connect with alumni for referrals.", "Apply for executive/analyst roles."]}
        ]
        
    # Sort roadmap by week names logically
    roadmap.sort(key=lambda x: x["week"])
    return weak_areas, roadmap

def get_company_recommendations(readiness_score: int, cgpa: float, programming_skills: int, technical_skills: int, internship: str, department: str, resume_keywords: list = None):
    if resume_keywords is None:
        resume_keywords = []
        
    resume_lower = [k.lower() for k in resume_keywords]
    
    # Map department keywords to canonical departments
    canonical_dept = department
    if any(keyword in department for keyword in ["AIML", "AIDS", "Artificial Intelligence", "Machine Learning", "Data Science"]):
        canonical_dept = "Artificial Intelligence & Machine Learning"
    
    # Standard inferred skills based on department
    dept_skills = []
    if canonical_dept == "Artificial Intelligence & Machine Learning":
        dept_skills = ["Python", "Machine Learning", "Deep Learning", "Data Science"]
    elif canonical_dept in ["Computer Science", "Information Technology"]:
        dept_skills = ["DSA", "Java", "Web Development"]
    elif canonical_dept == "Electronics & Communication":
        dept_skills = ["Embedded Systems", "IoT"]
    else:
        dept_skills = ["Aptitude", "Problem Solving", "Communication"]
        
    all_companies = [
        # Dream / Tier-1 Companies
        {
            "name": "NVIDIA",
            "skills": "Python, Machine Learning, Deep Learning, CUDA",
            "package": "₹18 - ₹36 LPA",
            "tier": "Dream",
            "req_cgpa": 8.2,
            "req_prog": 80,
            "req_tech": 80,
            "dept_boost": {"Artificial Intelligence & Machine Learning": 25, "Computer Science": 10},
            "core_skills": ["python", "machine learning", "deep learning", "data science"]
        },
        {
            "name": "Google",
            "skills": "Data Structures, Algorithms, Python/Java, System Design",
            "package": "₹15 - ₹35 LPA",
            "tier": "Dream",
            "req_cgpa": 8.0,
            "req_prog": 85,
            "req_tech": 80,
            "dept_boost": {"Computer Science": 20, "Artificial Intelligence & Machine Learning": 20, "Information Technology": 15},
            "core_skills": ["dsa", "java", "python", "machine learning"]
        },
        {
            "name": "Microsoft",
            "skills": "DSA, C++, Java, OOPs, Web Development",
            "package": "₹12 - ₹25 LPA",
            "tier": "Dream",
            "req_cgpa": 8.0,
            "req_prog": 80,
            "req_tech": 80,
            "dept_boost": {"Computer Science": 20, "Artificial Intelligence & Machine Learning": 20, "Information Technology": 15},
            "core_skills": ["dsa", "java", "web development"]
        },
        {
            "name": "Amazon",
            "skills": "DSA, Java, Web Development, System Design",
            "package": "₹14 - ₹28 LPA",
            "tier": "Dream",
            "req_cgpa": 7.5,
            "req_prog": 80,
            "req_tech": 80,
            "dept_boost": {"Computer Science": 20, "Artificial Intelligence & Machine Learning": 15, "Information Technology": 20},
            "core_skills": ["dsa", "java", "web development"]
        },
        {
            "name": "Intel",
            "skills": "C++, Verilog, Embedded Systems, Computer Architecture",
            "package": "₹14 - ₹26 LPA",
            "tier": "Dream",
            "req_cgpa": 8.0,
            "req_prog": 75,
            "req_tech": 80,
            "dept_boost": {"Electronics & Communication": 30, "Computer Science": 10},
            "core_skills": ["embedded systems", "iot"]
        },
        {
            "name": "Qualcomm",
            "skills": "C/C++, Embedded Systems, IoT, Digital Electronics",
            "package": "₹12 - ₹24 LPA",
            "tier": "Dream",
            "req_cgpa": 7.8,
            "req_prog": 75,
            "req_tech": 80,
            "dept_boost": {"Electronics & Communication": 30, "Computer Science": 10},
            "core_skills": ["embedded systems", "iot"]
        },
        {
            "name": "Texas Instruments",
            "skills": "Analog Electronics, Microcontrollers, Embedded C",
            "package": "₹12 - ₹22 LPA",
            "tier": "Dream",
            "req_cgpa": 7.8,
            "req_prog": 70,
            "req_tech": 80,
            "dept_boost": {"Electronics & Communication": 30},
            "core_skills": ["embedded systems", "iot"]
        },
        {
            "name": "Infosys AI Division",
            "skills": "Python, Machine Learning, Deep Learning, Data Science",
            "package": "₹6.5 - ₹10 LPA",
            "tier": "Dream",
            "req_cgpa": 7.0,
            "req_prog": 70,
            "req_tech": 70,
            "dept_boost": {"Artificial Intelligence & Machine Learning": 25, "Computer Science": 15, "Information Technology": 15},
            "core_skills": ["python", "machine learning", "data science"]
        },
        {
            "name": "Accenture",
            "skills": "Java, Web Development, DSA, Cloud Basics",
            "package": "₹4.5 - ₹10 LPA",
            "tier": "Dream",
            "req_cgpa": 6.5,
            "req_prog": 65,
            "req_tech": 65,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10},
            "core_skills": ["dsa", "java", "web development"]
        },
        
        # Service Tier Companies
        {
            "name": "TCS",
            "skills": "Aptitude, Java, Python, SQL, DSA Basics",
            "package": "₹3.6 - ₹7.5 LPA",
            "tier": "Service",
            "req_cgpa": 6.0,
            "req_prog": 55,
            "req_tech": 55,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10, "Mechanical Engineering": 5, "Civil Engineering": 5},
            "core_skills": ["dsa", "java", "web development", "python"]
        },
        {
            "name": "Infosys",
            "skills": "Aptitude, Python, Java, Web Development, DBMS",
            "package": "₹3.6 - ₹6.8 LPA",
            "tier": "Service",
            "req_cgpa": 6.0,
            "req_prog": 55,
            "req_tech": 55,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10, "Mechanical Engineering": 5, "Civil Engineering": 5},
            "core_skills": ["dsa", "java", "web development"]
        },
        {
            "name": "Wipro",
            "skills": "Coding, Java, DBMS, OOPs",
            "package": "₹3.5 - ₹6 LPA",
            "tier": "Service",
            "req_cgpa": 6.0,
            "req_prog": 55,
            "req_tech": 55,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10, "Mechanical Engineering": 5, "Civil Engineering": 5},
            "core_skills": ["dsa", "java"]
        },
        {
            "name": "Cognizant",
            "skills": "Programming basics, DBMS, SQL, Aptitude",
            "package": "₹4 - ₹6.5 LPA",
            "tier": "Service",
            "req_cgpa": 6.0,
            "req_prog": 55,
            "req_tech": 55,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10, "Mechanical Engineering": 5, "Civil Engineering": 5},
            "core_skills": ["dsa", "java", "web development"]
        },
        {
            "name": "Capgemini",
            "skills": "Pseudo code, Coding, Communication, SQL",
            "package": "₹4 - ₹7.5 LPA",
            "tier": "Service",
            "req_cgpa": 6.0,
            "req_prog": 55,
            "req_tech": 55,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10, "Mechanical Engineering": 5, "Civil Engineering": 5},
            "core_skills": ["dsa", "java", "web development"]
        },

        # Startup / Internships
        {
            "name": "Tech Startups (Zomato/Swiggy)",
            "skills": "Web Development, React, Node.js, Python",
            "package": "₹8 - ₹14 LPA",
            "tier": "Startup",
            "req_cgpa": 6.5,
            "req_prog": 70,
            "req_tech": 70,
            "dept_boost": {"Computer Science": 20, "Information Technology": 20, "Artificial Intelligence & Machine Learning": 15},
            "core_skills": ["web development", "python"]
        },
        {
            "name": "Internship programs (Hexaware/Virtusa)",
            "skills": "Basics of Programming, SQL, Aptitude",
            "package": "₹3.0 - ₹5.0 LPA",
            "tier": "Startup",
            "req_cgpa": 6.0,
            "req_prog": 50,
            "req_tech": 50,
            "dept_boost": {"Computer Science": 15, "Information Technology": 15, "Electronics & Communication": 10},
            "core_skills": ["dsa", "java", "web development"]
        },
        {
            "name": "Training-based companies (Revature/FDM)",
            "skills": "Basic Aptitude, Dedication, Learnability",
            "package": "₹3 - ₹4 LPA",
            "tier": "Startup",
            "req_cgpa": 5.5,
            "req_prog": 45,
            "req_tech": 45,
            "dept_boost": {"Computer Science": 10, "Information Technology": 10, "Electronics & Communication": 10, "Mechanical Engineering": 10, "Civil Engineering": 10},
            "core_skills": []
        }
    ]

    recommended = []
    
    for c in all_companies:
        # 1. Base score starts at 50
        match_score = 50
        
        # 2. Department Affinity Boost
        boost = c["dept_boost"].get(canonical_dept, 0)
        match_score += boost
        
        # If the department is completely unrelated to the company's focus, apply a penalty
        if boost == 0 and canonical_dept not in ["Computer Science", "Information Technology", "Artificial Intelligence & Machine Learning"] and c["tier"] == "Dream":
            match_score -= 25
            
        # 3. Academic & Technical Score Checks (Penalties)
        if cgpa < c["req_cgpa"]:
            match_score -= (c["req_cgpa"] - cgpa) * 15
        else:
            match_score += min((cgpa - c["req_cgpa"]) * 5, 10)
            
        if programming_skills < c["req_prog"]:
            match_score -= (c["req_prog"] - programming_skills) * 0.8
        else:
            match_score += min((programming_skills - c["req_prog"]) * 0.2, 5)
            
        if technical_skills < c["req_tech"]:
            match_score -= (c["req_tech"] - technical_skills) * 0.6
        else:
            match_score += min((technical_skills - c["req_tech"]) * 0.2, 5)
            
        # 4. Skills Match Boost
        # Check if company core skills overlap with department skills
        overlap = [s for s in c["core_skills"] if s.lower() in [ds.lower() for ds in dept_skills]]
        match_score += len(overlap) * 5
        
        # Check resume keywords
        if resume_lower:
            keyword_overlap = [k for k in c["core_skills"] if k.lower() in resume_lower]
            match_score += len(keyword_overlap) * 5
            
        # 5. Internship bonus
        if internship == "Yes" and c["tier"] in ["Dream", "Startup"]:
            match_score += 8
            
        # 6. Readiness score alignment
        # Higher readiness score should pull matches up, lower should pull down
        match_score += (readiness_score - 60) * 0.2
        
        # Limit matches to realistic range [30%, 98%]
        final_match = int(max(30, min(98, match_score)))
        
        recommended.append({
            "name": c["name"],
            "skills": c["skills"],
            "package": c["package"],
            "tier": c["tier"],
            "match": final_match
        })
        
    recommended.sort(key=lambda x: x["match"], reverse=True)
    
    # Filter recommendations: return only matches >= 65. If none, return top 3.
    filtered_recommended = [c for c in recommended if c["match"] >= 65]
    if not filtered_recommended:
        filtered_recommended = recommended[:3]
    
    if readiness_score >= 75 and cgpa >= 8.0:
        profile = "High Profile"
    elif readiness_score >= 50 or cgpa >= 7.0:
        profile = "Medium Profile"
    else:
        profile = "Beginner Profile"
        
    return profile, filtered_recommended

def run_predictions(p: StudentProfile):
    global model, salary_model
    
    # 1. Map to input features array for ML model
    dept_val = DEPT_MAP.get(p.department, 0)
    gender_val = 1 if p.gender == "Male" else 0
    internship_val = 1 if p.internship == "Yes" else 0
    resume_val = 1 if p.resume_uploaded == "Yes" else 0
    
    features = [
        dept_val,
        gender_val,
        p.tenth_percentage,
        p.twelfth_percentage,
        p.cgpa,
        p.backlogs,
        p.programming_skills,
        p.aptitude_score,
        p.communication_skills,
        p.technical_skills,
        p.projects,
        internship_val,
        p.certifications,
        p.hackathons,
        resume_val,
        p.mock_interview_score
    ]
    
    features_arr = np.array([features])
    
    # 2. Classifier Prediction
    placed_pred = 0
    prob_pred = 0.5
    
    if model is not None:
        try:
            placed_pred = int(model.predict(features_arr)[0])
            prob_pred = float(model.predict_proba(features_arr)[0][1])
        except Exception as e:
            print(f"Prediction failed with classifier: {e}. Falling back to rule-based.")
            # Fallback
            readiness_temp = calculate_readiness(p)
            prob_pred = readiness_temp / 100.0
    else:
        # Fallback
        readiness_temp = calculate_readiness(p)
        prob_pred = readiness_temp / 100.0
        
    # Smooth, realistic probability calibration
    # 1. CGPA adjustments
    if p.cgpa < 6.0:
        prob_pred -= 0.40
    elif p.cgpa < 6.5:
        prob_pred -= 0.20
    elif p.cgpa >= 8.5:
        prob_pred += 0.05
        
    # 2. Backlogs Penalty
    if p.backlogs >= 3:
        prob_pred -= 0.50
    elif p.backlogs == 2:
        prob_pred -= 0.30
    elif p.backlogs == 1:
        prob_pred -= 0.15
        
    # 3. Critical Technical Skill Check
    if p.programming_skills < 50:
        prob_pred -= 0.15
    if p.technical_skills < 50:
        prob_pred -= 0.15
        
    # 4. Critical Aptitude Check
    if p.aptitude_score < 50:
        prob_pred -= 0.10
        
    # 5. Elite Profile boost
    if p.cgpa >= 8.5 and p.programming_skills >= 80 and p.aptitude_score >= 80 and p.backlogs == 0:
        prob_pred = max(prob_pred, 0.88)
        
    # Bounds check
    prob_pred = max(0.02, min(prob_pred, 0.98))
    placed_pred = 1 if prob_pred >= 0.50 else 0
    
    # Calculate confidence score based on the probability of the predicted class
    confidence_score = int(prob_pred * 100) if placed_pred == 1 else int((1 - prob_pred) * 100)
    confidence_score = max(50, min(confidence_score, 100))

    # 3. Expected Salary package prediction
    salary_avg = 3.5
    if salary_model is not None:
        try:
            salary_avg = float(salary_model.predict(features_arr)[0])
        except Exception as e:
            print(f"Prediction failed with salary regressor: {e}.")
            # Fallback formula
            salary_avg = 3.2 + (p.cgpa - 5.0) * 1.5 + (p.programming_skills - 40.0) / 10.0 * 0.5 + p.projects * 0.8
    else:
        salary_avg = 3.2 + (p.cgpa - 5.0) * 1.5 + (p.programming_skills - 40.0) / 10.0 * 0.5 + p.projects * 0.8
        
    # Add bounds
    salary_avg = max(3.0, min(round(salary_avg, 2), 18.0))
    salary_low = max(3.0, round(salary_avg * 0.8, 2))
    salary_high = max(4.0, round(salary_avg * 1.3, 2))
    
    # Placed/Not placed check override: if not placed, package potential is lower but still show estimation
    
    # 4. Readiness score
    readiness_score = calculate_readiness(p)
    
    # 5. Reason for prediction
    reasons = []
    if p.cgpa >= 8.0:
        reasons.append("Excellent academic track record (CGPA)")
    elif p.cgpa < 6.5:
        reasons.append("Needs academic improvement (CGPA < 6.5)")
        
    if p.programming_skills >= 80 and p.technical_skills >= 80:
        reasons.append("Strong technical & programming skills")
    elif p.programming_skills < 60:
        reasons.append("Weak programming foundation")
        
    if p.internship == "Yes":
        reasons.append("Valuable hands-on internship experience")
    if p.projects >= 2:
        reasons.append(f"Completed {p.projects} relevant projects")
    if p.backlogs > 0:
        reasons.append(f"Active/history of backlogs ({p.backlogs})")
        
    if len(reasons) == 0:
        if placed_pred == 1:
            reasons.append("Balanced profile across technical, aptitude, and communication criteria")
        else:
            reasons.append("Requires improvement in multiple core skill domains")
            
    prediction_reason = " | ".join(reasons)
    
    # 6. Weak areas and Roadmap
    weak_areas, roadmap = get_weak_areas_and_roadmap(p)
    
    # 7. Recommendations
    _, recommended_companies = get_company_recommendations(
        readiness_score=readiness_score,
        cgpa=p.cgpa,
        programming_skills=p.programming_skills,
        technical_skills=p.technical_skills,
        internship=p.internship,
        department=p.department
    )
    
    return {
        "placed": placed_pred,
        "probability": prob_pred,
        "readiness_score": readiness_score,
        "salary_low": salary_low,
        "salary_avg": salary_avg,
        "salary_high": salary_high,
        "prediction_reason": prediction_reason,
        "weak_areas": weak_areas,
        "learning_roadmap": roadmap,
        "recommended_companies": recommended_companies,
        "programming_languages": p.programming_languages,
        "confidence_score": confidence_score
    }

# REST Endpoints
@app.post("/api/predict")
def predict_placement(p: StudentProfile):
    res = run_predictions(p)
    
    # Save/Update in SQLite database
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Check if register no already exists
    cursor.execute("SELECT id FROM students WHERE register_no = ?", (p.register_no,))
    existing = cursor.fetchone()
    
    # Create resume initial scores (default values unless updated by resume upload)
    res_score = 65
    missing_skills = ["data structures", "machine learning"]
    missing_keywords = ["Leadership", "Teamwork"]
    strengths = ["Technical Skills", "Academic Scores"]
    weaknesses = ["Missing leadership highlights", "Aptitude Score"]
    suggestions = ["Add communication examples", "Highlight group work"]
    
    if existing:
        # Keep old resume details if already there
        cursor.execute("""
            SELECT resume_score, resume_missing_skills, resume_missing_keywords, 
                   resume_strengths, resume_weaknesses, resume_suggestions 
            FROM students WHERE id = ?
        """, (existing["id"],))
        old_res = cursor.fetchone()
        if old_res and old_res["resume_score"] is not None:
            res_score = old_res["resume_score"]
            missing_skills = json.loads(old_res["resume_missing_skills"])
            missing_keywords = json.loads(old_res["resume_missing_keywords"])
            strengths = json.loads(old_res["resume_strengths"])
            weaknesses = json.loads(old_res["resume_weaknesses"])
            suggestions = json.loads(old_res["resume_suggestions"])
            
        cursor.execute("""
            UPDATE students SET
                student_name = ?, department = ?, gender = ?, tenth_percentage = ?,
                twelfth_percentage = ?, cgpa = ?, backlogs = ?, programming_skills = ?,
                aptitude_score = ?, communication_skills = ?, technical_skills = ?,
                projects = ?, internship = ?, certifications = ?, hackathons = ?,
                resume_uploaded = ?, mock_interview_score = ?, placed = ?, probability = ?,
                readiness_score = ?, salary_low = ?, salary_avg = ?, salary_high = ?,
                prediction_reason = ?, weak_areas = ?, learning_roadmap = ?, recommended_companies = ?,
                resume_score = ?, resume_missing_skills = ?, resume_missing_keywords = ?,
                resume_strengths = ?, resume_weaknesses = ?, resume_suggestions = ?,
                programming_languages = ?, confidence_score = ?
            WHERE id = ?
        """, (
            p.student_name, p.department, p.gender, p.tenth_percentage, p.twelfth_percentage,
            p.cgpa, p.backlogs, p.programming_skills, p.aptitude_score, p.communication_skills,
            p.technical_skills, p.projects, p.internship, p.certifications, p.hackathons,
            p.resume_uploaded, p.mock_interview_score, res["placed"], res["probability"],
            res["readiness_score"], res["salary_low"], res["salary_avg"], res["salary_high"],
            res["prediction_reason"], json.dumps(res["weak_areas"]), json.dumps(res["learning_roadmap"]),
            json.dumps(res["recommended_companies"]), res_score, json.dumps(missing_skills),
            json.dumps(missing_keywords), json.dumps(strengths), json.dumps(weaknesses),
            json.dumps(suggestions), p.programming_languages, res["confidence_score"], existing["id"]
        ))
        student_id = existing["id"]
    else:
        cursor.execute("""
            INSERT INTO students (
                student_name, register_no, department, gender, tenth_percentage,
                twelfth_percentage, cgpa, backlogs, programming_skills,
                aptitude_score, communication_skills, technical_skills,
                projects, internship, certifications, hackathons,
                resume_uploaded, mock_interview_score, placed, probability,
                readiness_score, salary_low, salary_avg, salary_high,
                prediction_reason, weak_areas, learning_roadmap, recommended_companies,
                resume_score, resume_missing_skills, resume_missing_keywords,
                resume_strengths, resume_weaknesses, resume_suggestions,
                programming_languages, confidence_score
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            p.student_name, p.register_no, p.department, p.gender, p.tenth_percentage, p.twelfth_percentage,
            p.cgpa, p.backlogs, p.programming_skills, p.aptitude_score, p.communication_skills,
            p.technical_skills, p.projects, p.internship, p.certifications, p.hackathons,
            p.resume_uploaded, p.mock_interview_score, res["placed"], res["probability"],
            res["readiness_score"], res["salary_low"], res["salary_avg"], res["salary_high"],
            res["prediction_reason"], json.dumps(res["weak_areas"]), json.dumps(res["learning_roadmap"]),
            json.dumps(res["recommended_companies"]), res_score, json.dumps(missing_skills),
            json.dumps(missing_keywords), json.dumps(strengths), json.dumps(weaknesses),
            json.dumps(suggestions), p.programming_languages, res["confidence_score"]
        ))
        student_id = cursor.lastrowid
        
    conn.commit()
    
    # Fetch complete student profile
    cursor.execute("SELECT * FROM students WHERE id = ?", (student_id,))
    row = cursor.fetchone()
    
    conn.close()
    
    student = dict(row)
    student["weak_areas"] = json.loads(student["weak_areas"]) if student.get("weak_areas") else []
    student["learning_roadmap"] = json.loads(student["learning_roadmap"]) if student.get("learning_roadmap") else []
    student["recommended_companies"] = json.loads(student["recommended_companies"]) if student.get("recommended_companies") else []
    student["resume_missing_skills"] = json.loads(student["resume_missing_skills"]) if student.get("resume_missing_skills") else []
    student["resume_missing_keywords"] = json.loads(student["resume_missing_keywords"]) if student.get("resume_missing_keywords") else []
    student["resume_strengths"] = json.loads(student["resume_strengths"]) if student.get("resume_strengths") else []
    student["resume_weaknesses"] = json.loads(student["resume_weaknesses"]) if student.get("resume_weaknesses") else []
    student["resume_suggestions"] = json.loads(student["resume_suggestions"]) if student.get("resume_suggestions") else []
    
    # Interview fields
    student["interview_questions"] = json.loads(student["interview_questions"]) if student.get("interview_questions") else []
    student["interview_responses"] = json.loads(student["interview_responses"]) if student.get("interview_responses") else {}
    student["interview_analysis"] = json.loads(student["interview_analysis"]) if student.get("interview_analysis") else {}
    student["interview_suggestions"] = json.loads(student["interview_suggestions"]) if student.get("interview_suggestions") else []
    
    return student

@app.get("/api/students")
def get_students(search: Optional[str] = None, dept: Optional[str] = None, sort: Optional[str] = None):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    query = "SELECT * FROM students WHERE 1=1"
    params = []
    
    if search:
        query += " AND (student_name LIKE ? OR register_no LIKE ?)"
        params.extend([f"%{search}%", f"%{search}%"])
        
    if dept and dept != "All":
        query += " AND department = ?"
        params.append(dept)
        
    if sort:
        if sort == "cgpa_desc":
            query += " ORDER BY cgpa DESC"
        elif sort == "cgpa_asc":
            query += " ORDER BY cgpa ASC"
        elif sort == "readiness_desc":
            query += " ORDER BY readiness_score DESC"
        elif sort == "readiness_asc":
            query += " ORDER BY readiness_score ASC"
        elif sort == "name_asc":
            query += " ORDER BY student_name ASC"
        else:
            query += " ORDER BY id DESC"
    else:
        query += " ORDER BY id DESC"
        
    cursor.execute(query, params)
    rows = cursor.fetchall()
    conn.close()
    
    students = []
    for r in rows:
        students.append(dict(r))
    return students

@app.get("/api/students/{register_no}")
def get_student_by_reg(register_no: str):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (register_no,))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        raise HTTPException(status_code=404, detail="Student profile not found")
        
    student = dict(row)
    student["weak_areas"] = json.loads(student["weak_areas"]) if student.get("weak_areas") else []
    student["learning_roadmap"] = json.loads(student["learning_roadmap"]) if student.get("learning_roadmap") else []
    student["recommended_companies"] = json.loads(student["recommended_companies"]) if student.get("recommended_companies") else []
    student["resume_missing_skills"] = json.loads(student["resume_missing_skills"]) if student.get("resume_missing_skills") else []
    student["resume_missing_keywords"] = json.loads(student["resume_missing_keywords"]) if student.get("resume_missing_keywords") else []
    student["resume_strengths"] = json.loads(student["resume_strengths"]) if student.get("resume_strengths") else []
    student["resume_weaknesses"] = json.loads(student["resume_weaknesses"]) if student.get("resume_weaknesses") else []
    student["resume_suggestions"] = json.loads(student["resume_suggestions"]) if student.get("resume_suggestions") else []
    
    # Interview fields
    student["interview_questions"] = json.loads(student["interview_questions"]) if student.get("interview_questions") else []
    student["interview_responses"] = json.loads(student["interview_responses"]) if student.get("interview_responses") else {}
    student["interview_analysis"] = json.loads(student["interview_analysis"]) if student.get("interview_analysis") else {}
    student["interview_suggestions"] = json.loads(student["interview_suggestions"]) if student.get("interview_suggestions") else []
    
    return student

@app.delete("/api/students/{register_no}")
def delete_student(register_no: str):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT id FROM students WHERE register_no = ?", (register_no,))
    row = cursor.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=404, detail="Student profile not found")
        
    cursor.execute("DELETE FROM students WHERE register_no = ?", (register_no,))
    conn.commit()
    conn.close()
    return {"message": f"Student with registration number {register_no} deleted successfully"}

# Resume Upload & NLP Analysis
@app.post("/api/analyze-resume")
async def analyze_resume(
    register_no: str = Form(...),
    file: UploadFile = File(...)
):
    # Determine the file type
    filename = file.filename.lower()
    ext = os.path.splitext(filename)[1]
    
    if ext not in [".pdf", ".docx", ".pptx"]:
        raise HTTPException(status_code=400, detail="Unsupported file format. Please upload PDF, DOCX, or PPTX.")
        
    temp_dir = "temp_uploads"
    os.makedirs(temp_dir, exist_ok=True)
    temp_filepath = os.path.join(temp_dir, file.filename)
    
    # Save the uploaded file locally
    with open(temp_filepath, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    # Parse text
    extracted_text = ""
    try:
        if ext == ".pdf":
            reader = PdfReader(temp_filepath)
            for page in reader.pages:
                text_page = page.extract_text()
                if text_page:
                    extracted_text += text_page + "\n"
        elif ext == ".docx":
            extracted_text = docx2txt.process(temp_filepath)
        elif ext == ".pptx":
            prs = Presentation(temp_filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + "\n"
    except Exception as e:
        if os.path.exists(temp_filepath):
            os.remove(temp_filepath)
        raise HTTPException(status_code=500, detail=f"Failed to parse resume: {e}")
        
    if os.path.exists(temp_filepath):
        os.remove(temp_filepath)
        
    # Analyze text with simple NLP / Keyword mapping
    text_lower = extracted_text.lower()
    
    found_keywords = []
    missing_keywords = []
    
    for kw in RESUME_KEYWORDS:
        if kw in text_lower:
            found_keywords.append(kw.title())
        else:
            missing_keywords.append(kw.title())
            
    # Calculate score based on keyword match density + general text length
    keyword_score = (len(found_keywords) / len(RESUME_KEYWORDS)) * 100
    
    # Formatting score
    formatting_score = 15
    if len(extracted_text) < 150:
        formatting_score = 5 # Too short/scanned pdf
        
    resume_score = int(min(keyword_score * 0.85 + formatting_score, 100))
    
    # Establish strengths and weaknesses
    strengths = []
    weaknesses = []
    suggestions = []
    missing_skills = []
    
    # Mock strengths and weaknesses
    if "python" in text_lower or "java" in text_lower:
        strengths.append("Demonstrated programming proficiency")
    if "projects" in text_lower:
        strengths.append("Project highlights present in resume")
    else:
        weaknesses.append("No project highlights found")
        suggestions.append("Add a detailed 'Projects' section describing 2-3 key technical projects with links.")
        
    if "internship" in text_lower:
        strengths.append("Professional industry/internship experience noted")
    else:
        weaknesses.append("Missing professional/internship experience")
        suggestions.append("Include internships, open-source contributions, or virtual experiences to showcase real-world work.")
        
    if "sql" not in text_lower:
        missing_skills.append("Database Management / SQL")
        suggestions.append("Add Database/SQL experience as most placement recruitment tests check for DBMS concepts.")
    if "machine learning" not in text_lower:
        missing_skills.append("Machine Learning / AI")
        
    if len(missing_keywords) > 0:
        weaknesses.append(f"Missing essential industry keywords: {', '.join(missing_keywords[:3])}")
        suggestions.append(f"Incorporate missing core keywords such as: {', '.join(missing_keywords)} naturally in your experience descriptions.")
        
    if resume_score >= 80:
        strengths.append("High resume text density and clean keyword targeting")
    elif resume_score < 60:
        weaknesses.append("Low keyword targeting score")
        suggestions.append("Improve your resume layout to ensure PDF parsers can easily read the text (avoid tables or complex icons).")
        
    if len(strengths) == 0:
        strengths.append("Clear structural headers")
    if len(weaknesses) == 0:
        weaknesses.append("Slightly brief profile detail")
        suggestions.append("Add a Professional Summary section at the top of your resume.")
        
    # Update SQLite database for this student profile
    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute("SELECT id FROM students WHERE register_no = ?", (register_no,))
    existing = cursor.fetchone()
    
    if existing:
        cursor.execute("""
            UPDATE students SET
                resume_uploaded = 'Yes',
                resume_score = ?,
                resume_missing_skills = ?,
                resume_missing_keywords = ?,
                resume_strengths = ?,
                resume_weaknesses = ?,
                resume_suggestions = ?
            WHERE register_no = ?
        """, (
            resume_score,
            json.dumps(missing_skills),
            json.dumps(missing_keywords),
            json.dumps(strengths),
            json.dumps(weaknesses),
            json.dumps(suggestions),
            register_no
        ))
        conn.commit()
    else:
        conn.close()
        raise HTTPException(status_code=404, detail="Student profile not found. Please create prediction profile first.")
        
    # Refetch updated student profile
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (register_no,))
    updated_row = cursor.fetchone()
    conn.close()
    
    return {
        "resume_score": resume_score,
        "missing_keywords": missing_keywords,
        "missing_skills": missing_skills,
        "strengths": strengths,
        "weaknesses": weaknesses,
        "suggestions": suggestions,
        "updated_profile": dict(updated_row) if updated_row else None
    }

class ChatMessage(BaseModel):
    message: str
    register_no: Optional[str] = None
    api_key: Optional[str] = None
    history: Optional[List[dict]] = None

def call_gemini(api_key: str, contents: list, system_instruction: str = None) -> str:
    import urllib.request
    import urllib.error
    import json
    import time
    
    # Try different models in case of failure or deprecation
    models = ["gemini-1.5-flash", "gemini-2.0-flash", "gemini-2.5-flash"]
    
    last_error = ""
    for model_name in models:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
        
        body = {
            "contents": contents
        }
        
        if system_instruction:
            body["system_instruction"] = {
                "parts": [{"text": system_instruction}]
            }
            
        headers = {
            "Content-Type": "application/json"
        }
        
        # Retry logic for rate limits and server errors
        max_retries = 3
        for attempt in range(max_retries):
            req = urllib.request.Request(
                url, 
                data=json.dumps(body).encode("utf-8"), 
                headers=headers, 
                method="POST"
            )
            try:
                with urllib.request.urlopen(req, timeout=12) as response:
                    res_data = json.loads(response.read().decode("utf-8"))
                    text = res_data["candidates"][0]["content"]["parts"][0]["text"]
                    return text
            except urllib.error.HTTPError as e:
                try:
                    err_msg = e.read().decode("utf-8")
                except Exception:
                    err_msg = "Unknown error"
                print(f"Gemini API HTTP Error ({model_name}, attempt {attempt+1}/{max_retries}): {e.code} - {err_msg}")
                last_error = f"HTTP {e.code}: {err_msg}"
                
                # If rate limit (429) or server error (5xx), wait and retry
                if e.code in [429, 500, 502, 503, 504]:
                    time.sleep(2 ** (attempt + 1))
                    continue
                else:
                    # For auth errors (400, 403), do not retry the same model
                    break
            except Exception as e:
                print(f"Gemini API Connection Error ({model_name}, attempt {attempt+1}/{max_retries}): {e}")
                last_error = str(e)
                time.sleep(2 ** (attempt + 1))
                continue
                
    return f"Error executing Gemini API: {last_error}"

@app.post("/api/chat")
def chat_guidance(msg: ChatMessage):
    user_query = msg.message.lower()
    
    # 1. Fetch student context if register number is provided
    student_details = None
    if msg.register_no:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM students WHERE register_no = ?", (msg.register_no,))
        row = cursor.fetchone()
        conn.close()
        if row:
            student_details = dict(row)
            
    # Resolve API Key (payload > environment > default)
    gemini_key = msg.api_key or os.getenv("GEMINI_API_KEY") or "AQ.Ab8RN6JIB7yz38nl6_zsdY3c5S2TzkbHLDui0N0F2C9VUFmHng"
    
    if student_details:
        name = student_details["student_name"]
        cgpa = student_details["cgpa"]
        readiness = student_details["readiness_score"]
        prob = int(student_details["probability"] * 100)
        dept = student_details["department"]
        weak_list = json.loads(student_details["weak_areas"])
        weak_areas_str = ", ".join([w["area"] for w in weak_list]) if weak_list else "None"
        res_score = student_details["resume_score"]
        mock_score = student_details["mock_interview_score"]
    else:
        name = "Student"
        cgpa = "N/A"
        readiness = "N/A"
        prob = "N/A"
        dept = "N/A"
        weak_areas_str = "N/A"
        res_score = "N/A"
        mock_score = "N/A"

    # If Gemini API key is available, attempt real-time GenAI call
    if gemini_key and gemini_key.strip():
        system_instruction = (
            "You are a professional AI Career Guidance Coach. "
            f"You are speaking with {name}, a student from the {dept} department. "
            f"Academic/Skill context: CGPA is {cgpa}/10.0, Placement Readiness Score is {readiness}/100, "
            f"Placement Probability is {prob}%, Resume Score is {res_score}/100, Mock Interview Score is {mock_score}/100, "
            f"and weak areas are [{weak_areas_str}]. "
            "Respond to their career, resume, coding, or interview preparation questions. "
            "Give direct, helpful, and highly actionable answers. Avoid long preambles, and be supportive."
        )
        
        # Clean and construct contents from history
        contents = []
        for h in msg.history or []:
            role = h.get("role", "user")
            if role in ["bot", "model"]:
                role = "model"
            else:
                role = "user"
            
            text = h.get("text", "").strip()
            if not text:
                continue
                
            # First turn in Gemini conversation must be 'user'
            if not contents:
                if role == "model":
                    continue # Skip initial model messages
                contents.append({"role": role, "parts": [{"text": text}]})
            else:
                # Alternate roles strictly (combine consecutive same roles)
                if contents[-1]["role"] == role:
                    contents[-1]["parts"][0]["text"] += "\n" + text
                else:
                    contents.append({"role": role, "parts": [{"text": text}]})
        
        # Ensure current message alternates with the last turn
        if not contents or contents[-1]["role"] == "model":
            contents.append({
                "role": "user",
                "parts": [{"text": msg.message}]
            })
        else:
            # If last turn was user, append the new query text to it
            contents[-1]["parts"][0]["text"] += "\n" + msg.message
        
        try:
            reply = call_gemini(gemini_key.strip(), contents, system_instruction)
            if not reply.startswith("Error executing Gemini API"):
                return {"response": reply}
            else:
                print(f"Gemini API failure: {reply}. Switching to fallback career guidance mode.")
        except Exception as e:
            print(f"Gemini API exception: {e}. Switching to fallback career guidance mode.")
        
    # Generate bot response based on keywords (Local NLP Router Fallback)
    prev_user_msg = ""
    prev_bot_msg = ""
    if msg.history and len(msg.history) >= 2:
        for turn in reversed(msg.history):
            if turn.get("role") == "user" and not prev_user_msg:
                prev_user_msg = turn.get("text", "").lower()
            elif (turn.get("role") == "model" or turn.get("role") == "bot") and not prev_bot_msg:
                prev_bot_msg = turn.get("text", "").lower()

    is_follow_up = any(kw in user_query for kw in ["more", "detail", "elaborate", "how", "why", "explain", "projects", "example", "suggest"])
    
    # 1. Communication / Soft skills check
    if "communication" in user_query or "speaking" in user_query or "english" in user_query or "soft skill" in user_query or "interview communication" in user_query or (is_follow_up and "communication" in prev_user_msg):
        if student_details and student_details.get("communication_skills", 100) < 70:
            comm_score = student_details["communication_skills"]
            resp = (
                f"Hello {name}! I noticed your Communication Score is **{comm_score}/100**, which is below the target benchmark of 75%. "
                "Here is your personalized **30-Day Communication Improvement Plan**:\n\n"
                "1. **Daily Mirror Practice (15 mins)**: Answer standard HR questions (like 'Tell me about yourself') out loud in front of a mirror or record it on video. Review your posture and speed.\n"
                "2. **STAR Method Drafting**: Write answers for 5 behavioral questions using: **Situation, Task, Action, and Result**. This keeps your answers structured and concise.\n"
                "3. **Peer Mock Group Discussions**: Conduct weekly GDs with 3-4 classmates on current tech trends in English. Ask for honest feedback on your tone and articulation.\n"
                "4. **Active Listening & Vocabulary**: Read one technical article and summarize it in your own words. Practice talking slowly (aim for 120-140 words per minute).\n\n"
                "Would you like advice on answering standard HR behavioral questions?"
            )
        else:
            resp = (
                "For professional communication in placement rounds, practice these core techniques:\n"
                "- **The STAR Method**: Describe the **Situation**, your **Task**, the **Action** you took, and the **Result** achieved. Use this for all project and behavioral questions.\n"
                "- **Thinking Out Loud**: During technical rounds, communicate your thought process while coding. Interviewers value *how* you solve a problem as much as the solution itself.\n"
                "- **Elevator Pitch**: Prepare a crisp, 90-second self-introduction highlighting your top project, key skills, and career objective."
            )

    # 2. AIML Career & Skills check
    elif "aiml" in user_query or "artificial intelligence" in user_query or "machine learning" in user_query or "data science" in user_query or "ml" in user_query or (is_follow_up and ("aiml" in prev_user_msg or "machine learning" in prev_user_msg or dept == "Artificial Intelligence & Machine Learning")):
        if student_details and (student_details["department"] == "Artificial Intelligence & Machine Learning" or "aiml" in user_query or "ml" in user_query):
            resp = (
                f"Hi {name}! Since you are interested in AI/ML, here is a highly tailored career roadmap:\n\n"
                "**Essential Skills to Master**:\n"
                "- *Mathematics*: Linear Algebra, Calculus, and Applied Statistics (Probability, Hypothesis testing).\n"
                "- *Languages & Libraries*: Python, NumPy, Pandas, Scikit-Learn (regression, classification, clustering).\n"
                "- *Data Engineering*: SQL (joins, aggregation, subqueries), data wrangling, and feature scaling.\n"
                "- *Deep Learning*: Neural networks, CNNs, RNNs, and framework libraries like PyTorch or TensorFlow.\n\n"
                "**Target Career Paths**:\n"
                "1. **Machine Learning Engineer**: Focuses on deploying ML models, model monitoring, and pipeline engineering (MLOps).\n"
                "2. **Data Scientist**: Analyzes business data, builds statistical models, and delivers predictive insights.\n"
                "3. **AI Software Engineer**: Integrates pre-trained models (LLMs, computer vision models) into standard web applications.\n\n"
                "**Recommended Projects**:\n"
                "- *Beginner*: Build an end-to-end ML classification app (like this Placement Predictor!) using FastAPI and React.\n"
                "- *Advanced*: Train a customized fine-tuned model for text summarization or object detection, hosting it on HuggingFace/AWS."
            )
        else:
            resp = (
                "AI/ML is a high-growth domain. To start, learn Python programming and core mathematics (statistics & linear algebra). "
                "Next, practice data manipulation using Pandas and Scikit-Learn. Build regression and classification projects on Kaggle datasets, "
                "and learn how to deploy them using simple API servers like Flask or FastAPI."
            )

    # 3. ECE Career & Skills check
    elif "ece" in user_query or "electronics" in user_query or "communication engineering" in user_query or "embedded" in user_query or "iot" in user_query or "hardware" in user_query or (is_follow_up and ("ece" in prev_user_msg or "embedded" in prev_user_msg or "hardware" in prev_user_msg or dept == "Electronics & Communication")):
        if student_details and (student_details["department"] == "Electronics & Communication" or "ece" in user_query or "embedded" in user_query):
            resp = (
                f"Hi {name}! As an Electronics & Communication student, you have a unique advantage in both hardware and software. Here is your personalized core/embedded roadmap:\n\n"
                "**1. Core Embedded & IoT Skills**:\n"
                "- Master Embedded C and C++ programming.\n"
                "- Core subjects: Microcontrollers (8051, ARM Cortex, Arduino, Raspberry Pi), Digital Electronics, and Computer Architecture.\n"
                "- Learn RTOS (Real-Time Operating Systems) and communication protocols (UART, SPI, I2C, CAN).\n\n"
                "**2. Target Companies**:\n"
                "- **Tier-1 Core Hardware**: Intel, Qualcomm, Texas Instruments, NVIDIA, ARM.\n"
                "- **Systems & IoT**: Bosch, Samsung, Honeywell.\n\n"
                "**3. Recommended Projects**:\n"
                "- Build an IoT-based sensor monitoring system (using ESP32 and MQTT).\n"
                "- Design a custom microcontroller circuit or simulate digital logic systems on FPGA.\n\n"
                "Would you like recommendations on core ECE versus software development career paths?"
            )
        else:
            resp = (
                "For Electronics & Communication Engineering (ECE) roles, focus on core skills like Embedded Systems, "
                "microcontrollers, and digital electronics. Learn C/C++ programming and hardware descriptive languages (Verilog/VHDL). "
                "Target top core companies like Qualcomm, Intel, and Texas Instruments."
            )

    # 4. CSE Career & Skills check
    elif "cse" in user_query or "computer science" in user_query or "software engineer" in user_query or "software job" in user_query or "dsa" in user_query or "programming" in user_query or (is_follow_up and ("cse" in prev_user_msg or "software" in prev_user_msg or "coding" in prev_user_msg)):
        if student_details and (student_details["department"] == "Computer Science" or "cse" in user_query or "software" in user_query):
            resp = (
                f"Hi {name}! As a Computer Science student, your placement preparation should target core engineering fundamentals. Here is your dashboard-customized roadmap:\n\n"
                "**1. Core Coding & Data Structures (DSA)**:\n"
                "- Master one language: Java, C++, or Python.\n"
                "- Focus Areas: Arrays, HashMaps, Two-Pointers, Recursion, Trees, Graphs, and Dynamic Programming.\n"
                "- Target: Solve 150-200 curated questions on LeetCode (e.g., NeetCode 150 track).\n\n"
                "**2. CS Foundations**:\n"
                "- *DBMS*: Relational algebra, SQL queries (joins, indexing), ACID properties, and normalization.\n"
                "- *Operating Systems*: Process management, multithreading, deadlocks, and virtual memory allocation.\n"
                "- *Computer Networks*: TCP/IP model, HTTP/S headers, DNS, and socket connections.\n\n"
                "**3. System Architecture & Projects**:\n"
                "- Develop a robust full-stack project (e.g., React front-end, SQLite/PostgreSQL database, FastAPI back-end).\n"
                "- Revise basic System Design: horizontal/vertical scaling, caching (Redis), and load balancers.\n\n"
                "**Recommended Roles**: Software Development Engineer (SDE), Backend Developer, DevOps Engineer, Site Reliability Engineer (SRE)."
            )
        else:
            resp = (
                "For Computer Science SDE roles, recruiters prioritize deep problem-solving skills. Focus heavily on mastering "
                "Data Structures & Algorithms (Trees, Graphs, DP) in C++, Java, or Python. Revise DBMS (SQL Joins), Object-Oriented "
                "Programming (OOPs), and Operating Systems. Host 2 detailed projects on GitHub demonstrating clean API architecture."
            )

    # 5. Resume check
    elif "resume" in user_query or "cv" in user_query or (is_follow_up and "resume" in prev_user_msg):
        if student_details and res_score is not None:
            miss_kw = json.loads(student_details["resume_missing_keywords"])
            miss_skills = json.loads(student_details["resume_missing_skills"])
            resp = (
                f"Hello {name}! Your current Resume NLP score is **{res_score}/100**.\n\n"
                f"**Missing Keywords Identified**: {', '.join(miss_kw[:4]) if miss_kw else 'None'}\n"
                f"**Missing Skill Areas**: {', '.join(miss_skills[:3]) if miss_skills else 'None'}\n\n"
                "**SaaS Optimization Tips**:\n"
                "- **ATS Formatting**: Avoid multi-column layouts, tables, and complex graphic progress bars. Use a clean, single-column design (like the Jake's Resume LaTeX template).\n"
                "- **Impact Metrics**: Use the Google XYZ formula: *'Accomplished [X], as measured by [Y], by doing [Z]'* (e.g., 'Optimized query latency by 25% by implementing Redis caching').\n"
                "- **Skills Section**: Group skills categorically (Languages, Frameworks, Databases, Tools) at the top of the resume."
            )
        else:
            resp = (
                "To optimize your resume for applicant tracking systems (ATS):\n"
                "1. Structure it in a single column using simple headings (Education, Experience, Projects, Skills).\n"
                "2. Categorize your technical skills: Languages, Databases, Frameworks, Developer Tools.\n"
                "3. Describe projects with clear links to GitHub and focus on active verbs: 'Developed', 'Optimized', 'Integrated'."
            )

    # 6. Interview check
    elif "interview" in user_query or "prepare" in user_query or "mock" in user_query or (is_follow_up and "interview" in prev_user_msg):
        if student_details:
            resp = (
                f"Hi {name}! Your current Mock Interview score is **{mock_score}/100**.\n"
                "To prepare for upcoming placement interviews, follow this structured plan:\n\n"
                "1. **Technical Interview (Round 1 & 2)**:\n"
                "   - Expect 1-2 coding questions. Communicate your thoughts continuously. Speak out loud as you trace the code.\n"
                "   - Expect questions on SQL (be ready to write Joins and Subqueries on the fly) and OOP concepts.\n"
                "2. **System Design (Premium Packages - ₹10+ LPA)**:\n"
                "   - Review database selection (SQL vs NoSQL), caching layers, and API design principles.\n"
                "3. **HR & Behavioral Round**:\n"
                "   - Prepare answers for: 'Tell me about yourself', 'What is your biggest failure?', 'Why this company?'\n"
                "   - Ensure all answers map back to your active projects and internship achievements."
            )
        else:
            resp = (
                "To prepare for interviews:\n"
                "1. **Self-Introduction**: Prepare a solid 2-minute summary of your projects and achievements.\n"
                "2. **Coding Round**: Practice writing code on a whiteboard or online compiler under time pressure.\n"
                "3. **CS Fundamentals**: Review DBMS, SQL, and OOPs concepts thoroughly.\n"
                "4. **Behavioral Questions**: Practice standard HR questions like 'Why do you want to join us?' or 'Tell me about a time you failed.'"
            )

    # 7. Salary check
    elif "salary" in user_query or "package" in user_query or "lpa" in user_query or "earn" in user_query or (is_follow_up and "salary" in prev_user_msg):
        if student_details:
            resp = (
                f"Hello {name}, based on our ML Regressor evaluation, your expected starting salary range is **₹{student_details['salary_low']} LPA - ₹{student_details['salary_high']} LPA** "
                f"with a median prediction of **₹{student_details['salary_avg']} LPA**.\n\n"
                "**Key Salary Factors**:\n"
                "- *Academics*: Maintaining a CGPA above 8.0 qualifies you for dream tier brackets.\n"
                "- *Programming Score*: Reaching 80%+ on your programming skills unlocks premium developer packages (₹10+ LPA).\n"
                "- *Internship*: Having industry experience adds a 15% negotiation leverage to your baseline offer.\n\n"
                "Practice advanced DSA and system design to target the high end of your salary prediction!"
            )
        else:
            resp = (
                "The Placement System predicts starting salary packages based on student credentials. "
                "- Dream Tier (Google, Microsoft, Amazon): ₹12 - ₹35 LPA. Requires CGPA > 8.0 and coding score > 80.\n"
                "- Corporate Service Tier (TCS, Infosys, Wipro): ₹3.5 - ₹7 LPA. Requires CGPA > 6.0.\n"
                "- Startups / Training: ₹3.0 - ₹5.0 LPA."
            )

    # 8. Company check
    elif "company" in user_query or "companies" in user_query or "suitable" in user_query or "recommend" in user_query or (is_follow_up and "company" in prev_user_msg):
        if student_details:
            companies = json.loads(student_details["recommended_companies"])
            c_names = ", ".join([c["name"] for c in companies[:3]])
            resp = (
                f"Based on your placement readiness score (**{readiness}/100**), "
                f"the most suitable companies for your profile are: **{c_names}**.\n\n"
                "Here is the breakdown:\n" +
                "\n".join([f"- **{c['name']}** (Expected Package: {c['package']} | Match: {c['match']}%)" for c in companies[:4]]) +
                "\n\nYou can review all recommended companies in the 'Company Recommender' tab!"
            )
        else:
            resp = (
                "Based on student scores, our recommender categorizes suitability into three tiers:\n"
                "- **Dream Tier (Tier-1)**: Google, Microsoft, Amazon, Accenture Premium (₹10+ LPA).\n"
                "- **Service Tier (Tier-2)**: TCS, Infosys, Wipro, Cognizant (₹3.5 - ₹7 LPA).\n"
                "- **Startups**: Fast-paced tech startups hiring frontend/backend developers."
            )

    # 9. General placement / improvement check
    elif "improve" in user_query or "chance" in user_query or "readiness" in user_query or "placement" in user_query or "eligible" in user_query:
        if student_details:
            if weak_list:
                areas_str = ", ".join([w["area"] for w in weak_list])
                resp = (
                    f"Hello {name}! Based on your profile, you have a placement readiness score of **{readiness}/100** "
                    f"and a probability of **{prob}%** of getting placed.\n\n"
                    f"To improve your chances, focus on these weak areas: **{areas_str}**.\n"
                    f"Specifically, we recommend:\n" + 
                    "\n".join([f"- *{w['area']}*: {w['suggestion']}" for w in weak_list[:3]])
                )
            else:
                resp = (
                    f"Great job, {name}! Your readiness score is a high **{readiness}/100**. "
                    f"You have a **{prob}%** placement probability. We recommend practicing company-specific "
                    f"mock interviews (Google/Microsoft/Amazon) and revising system design questions to land premium packages."
                )
        else:
            resp = (
                "To improve your placement chances: \n"
                "1. Maintain a CGPA above 8.0 to clear company cutoffs.\n"
                "2. Solve 2 problems daily on LeetCode/GeeksforGeeks to strengthen DSA.\n"
                "3. Work on at least 2 major coding projects using modern frameworks (React, FastAPI, Node.js).\n"
                "4. Practice mock interviews to boost communication and technical explanation skills."
            )

    # 10. Default coach menu
    else:
        resp = (
            f"Hello {name}! I am your AI Career Guidance Assistant. I can help you with:\n"
            "- Suggestions to improve your placement readiness score.\n"
            "- Personalized skill advice and roadmaps.\n"
            "- Tips to optimize your resume score.\n"
            "- Company recommendations based on your scores.\n\n"
            "Ask me anything like: *'How to improve my resume?'* or *'Which companies are suitable for me?'*\n"
            "*(💡 Tip: Enter a Gemini API Key in Chat Settings to enable real-time dynamic answers on any topic!)*"
        )
        
    return {"response": resp}

# Admin Analytics API
@app.get("/api/analytics")
def get_analytics():
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # 1. General Stats
    cursor.execute("SELECT COUNT(*) FROM students")
    total_students = cursor.fetchone()[0]
    
    if total_students == 0:
        conn.close()
        return {
            "total_students": 0,
            "placed_students": 0,
            "not_placed_students": 0,
            "avg_readiness": 0,
            "dept_analytics": [],
            "ml_models": []
        }
        
    cursor.execute("SELECT COUNT(*) FROM students WHERE placed = 1")
    placed_students = cursor.fetchone()[0]
    
    not_placed_students = total_students - placed_students
    
    cursor.execute("SELECT AVG(readiness_score) FROM students")
    avg_readiness = round(cursor.fetchone()[0], 2)
    
    # 2. Department-wise analytics
    cursor.execute("""
        SELECT department, 
               COUNT(*) as total, 
               SUM(CASE WHEN placed = 1 THEN 1 ELSE 0 END) as placed_count,
               AVG(readiness_score) as avg_readiness
        FROM students
        GROUP BY department
    """)
    dept_rows = cursor.fetchall()
    
    dept_analytics = []
    for row in dept_rows:
        dept_analytics.append({
            "department": row["department"],
            "total": row["total"],
            "placed": row["placed_count"],
            "placement_rate": round((row["placed_count"] / row["total"]) * 100, 2),
            "avg_readiness": round(row["avg_readiness"], 2)
        })
        
    # 3. Model metrics comparison
    cursor.execute("SELECT * FROM ml_metrics")
    ml_rows = cursor.fetchall()
    
    ml_models = []
    for row in ml_rows:
        ml_models.append({
            "model_name": row["model_name"],
            "accuracy": row["accuracy"],
            "precision": row["precision"],
            "recall": row["recall"],
            "f1_score": row["f1_score"],
            "confusion_matrix": json.loads(row["confusion_matrix"]),
            "classification_report": json.loads(row["classification_report"]),
            "is_best": bool(row["is_best"])
        })
        
    conn.close()
    
    return {
        "total_students": total_students,
        "placed_students": placed_students,
        "not_placed_students": not_placed_students,
        "avg_readiness": avg_readiness,
        "dept_analytics": dept_analytics,
        "ml_models": ml_models
    }

# Seed DB Route (useful to test admin dashboard immediately with sample data)
@app.get("/api/seed-db")
def seed_database():
    csv_path = "dataset/placement_data.csv"
    if not os.path.exists(csv_path):
        raise HTTPException(status_code=404, detail="Dataset CSV not found. Run train_models.py first to generate CSV and train models.")
        
    df = pd.read_csv(csv_path)
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Clear existing entries
    cursor.execute("DELETE FROM students")
    
    inserted_count = 0
    # Seed top 100 students for quick loading
    for idx, row in df.head(100).iterrows():
        p = StudentProfile(
            student_name=row["student_name"],
            register_no=row["register_no"],
            department=row["department"],
            gender=row["gender"],
            tenth_percentage=float(row["tenth_percentage"]),
            twelfth_percentage=float(row["twelfth_percentage"]),
            cgpa=float(row["cgpa"]),
            backlogs=int(row["backlogs"]),
            programming_skills=int(row["programming_skills"]),
            aptitude_score=int(row["aptitude_score"]),
            communication_skills=int(row["communication_skills"]),
            technical_skills=int(row["technical_skills"]),
            projects=int(row["projects"]),
            internship=row["internship"],
            certifications=int(row["certifications"]),
            hackathons=int(row["hackathons"]),
            resume_uploaded=row["resume_uploaded"],
            mock_interview_score=int(row["mock_interview_score"])
        )
        
        # Calculate ML results
        res = run_predictions(p)
        
        # Resume mockup values for seeded records
        res_score = int(60 + (p.programming_skills * 0.2) + (p.communication_skills * 0.1) + min(p.projects * 5, 10))
        res_score = min(res_score, 100)
        
        missing_skills = []
        if p.programming_skills < 70:
            missing_skills.append("Advanced Coding")
        if p.technical_skills < 70:
            missing_skills.append("SQL / Databases")
            
        missing_keywords = []
        if p.projects == 0:
            missing_keywords.append("Projects")
        if p.internship == "No":
            missing_keywords.append("Internship")
        if p.communication_skills < 75:
            missing_keywords.append("Communication")
            
        strengths = ["Academic Consistency"]
        if p.cgpa >= 8.5:
            strengths.append("High GPA honors")
        if p.programming_skills >= 85:
            strengths.append("Strong Programming")
            
        weaknesses = []
        if p.backlogs > 0:
            weaknesses.append("Active backlogs")
        if p.mock_interview_score < 70:
            weaknesses.append("Weak interview communication")
            
        suggestions = ["Ensure clean formatting", "Add Github URLs for projects"]
        
        prog_langs = "Python, Java, SQL" if p.department in ["Computer Science", "Information Technology"] else ("Python, R, SQL" if p.department == "Artificial Intelligence & Machine Learning" else "C, C++, Embedded C")
        
        cursor.execute("""
            INSERT INTO students (
                student_name, register_no, department, gender, tenth_percentage,
                twelfth_percentage, cgpa, backlogs, programming_skills,
                aptitude_score, communication_skills, technical_skills,
                projects, internship, certifications, hackathons,
                resume_uploaded, mock_interview_score, placed, probability,
                readiness_score, salary_low, salary_avg, salary_high,
                prediction_reason, weak_areas, learning_roadmap, recommended_companies,
                resume_score, resume_missing_skills, resume_missing_keywords,
                resume_strengths, resume_weaknesses, resume_suggestions,
                programming_languages, confidence_score
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            p.student_name, p.register_no, p.department, p.gender, p.tenth_percentage, p.twelfth_percentage,
            p.cgpa, p.backlogs, p.programming_skills, p.aptitude_score, p.communication_skills,
            p.technical_skills, p.projects, p.internship, p.certifications, p.hackathons,
            p.resume_uploaded, p.mock_interview_score, res["placed"], res["probability"],
            res["readiness_score"], res["salary_low"], res["salary_avg"], res["salary_high"],
            res["prediction_reason"], json.dumps(res["weak_areas"]), json.dumps(res["learning_roadmap"]),
            json.dumps(res["recommended_companies"]), res_score, json.dumps(missing_skills),
            json.dumps(missing_keywords), json.dumps(strengths), json.dumps(weaknesses),
            json.dumps(suggestions), prog_langs, res["confidence_score"]
        ))
        inserted_count += 1
        
    conn.commit()
    conn.close()
    
    return {"message": f"Successfully seeded database with {inserted_count} student profiles from CSV."}

# Helper to generate PDF in-memory using reportlab
def generate_pdf_report(student: dict) -> BytesIO:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=54,
        bottomMargin=45
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#6d28d9'), # Purple accent
        alignment=0,
        spaceAfter=15
    )
    
    h1_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=14,
        textColor=colors.HexColor('#0f172a'),
        spaceBefore=10,
        spaceAfter=5,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=8.5,
        leading=11,
        textColor=colors.HexColor('#334155'),
        spaceAfter=3
    )
    
    bold_body_style = ParagraphStyle(
        'BoldBodyTextCustom',
        parent=body_style,
        fontName='Helvetica-Bold'
    )
    
    table_header_style = ParagraphStyle(
        'TableHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=8.5,
        leading=10,
        textColor=colors.white,
        alignment=0
    )
    
    table_cell_style = ParagraphStyle(
        'TableCell',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=7.5,
        leading=9.5,
        textColor=colors.HexColor('#1e293b')
    )
    
    story = []
    
    # Header Banner
    story.append(Paragraph("AI-POWERED CAREER INTELLIGENCE PLATFORM", ParagraphStyle('Sub', fontName='Helvetica-Bold', fontSize=8, leading=10, textColor=colors.HexColor('#94a3b8'), spaceAfter=2)))
    story.append(Paragraph("COMPREHENSIVE PERFORMANCE & PLACEMENT READINESS REPORT", title_style))
    story.append(Spacer(1, 8))
    
    # Section 1: Academic & Demographic Profile
    story.append(Paragraph("1. STUDENT PROFILE & ACADEMIC CREDENTIALS", h1_style))
    
    profile_data = [
        [Paragraph("Full Name", bold_body_style), Paragraph(student["student_name"], body_style),
         Paragraph("Register Number", bold_body_style), Paragraph(student["register_no"], body_style)],
        [Paragraph("Department", bold_body_style), Paragraph(student["department"], body_style),
         Paragraph("Gender", bold_body_style), Paragraph(student["gender"], body_style)],
        [Paragraph("10th Percentage", bold_body_style), Paragraph(f"{student['tenth_percentage']}%", body_style),
         Paragraph("12th Percentage", bold_body_style), Paragraph(f"{student['twelfth_percentage']}%", body_style)],
        [Paragraph("Current CGPA", bold_body_style), Paragraph(f"{student['cgpa']} / 10.0", bold_body_style),
         Paragraph("Active Backlogs", bold_body_style), Paragraph(str(student['backlogs']), body_style)],
        [Paragraph("Programming Score", bold_body_style), Paragraph(f"{student['programming_skills']}/100", body_style),
         Paragraph("Aptitude Score", bold_body_style), Paragraph(f"{student['aptitude_score']}/100", body_style)],
        [Paragraph("Communication Score", bold_body_style), Paragraph(f"{student['communication_skills']}/100", body_style),
         Paragraph("Technical Score", bold_body_style), Paragraph(f"{student['technical_skills']}/100", body_style)],
        [Paragraph("Projects Completed", bold_body_style), Paragraph(str(student['projects']), body_style),
         Paragraph("Internship", bold_body_style), Paragraph(student['internship'], body_style)],
        [Paragraph("Certifications", bold_body_style), Paragraph(str(student['certifications']), body_style),
         Paragraph("Coding Languages", bold_body_style), Paragraph(student.get('programming_languages', 'Python, Java, SQL'), body_style)]
    ]
    
    t1 = Table(profile_data, colWidths=[110, 160, 110, 160])
    t1.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('PADDING', (0,0), (-1,-1), 4),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ]))
    story.append(t1)
    story.append(Spacer(1, 10))
    
    # Section 2: Predictive Placement Metrics
    story.append(Paragraph("2. ML PREDICTIVE PLACEMENT EVALUATION", h1_style))
    
    placed_text = "ELIGIBLE FOR PLACEMENT (PLACED)" if student["placed"] == 1 else "RE-TRAINING & IMPROVEMENT REQUIRED (NOT PLACED)"
    placed_color = colors.HexColor('#059669') if student["placed"] == 1 else colors.HexColor('#dc2626')
    probability = int(student["probability"] * 100)
    confidence = student.get("confidence_score", int(student["probability"] * 100) if student["placed"] == 1 else int((1 - student["probability"]) * 100))
    readiness = student["readiness_score"]
    
    readiness_cat = "Needs Improvement"
    if readiness >= 90:
        readiness_cat = "Excellent"
    elif readiness >= 75:
        readiness_cat = "Good"
    elif readiness >= 60:
        readiness_cat = "Average"
        
    metrics_data = [
        [Paragraph("Placement Status Prediction", bold_body_style), Paragraph(placed_text, ParagraphStyle('PlColor', parent=bold_body_style, textColor=placed_color))],
        [Paragraph("Placement Probability", bold_body_style), Paragraph(f"{probability}%", body_style)],
        [Paragraph("Prediction Confidence Score", bold_body_style), Paragraph(f"{confidence}%", body_style)],
        [Paragraph("Placement Readiness Score", bold_body_style), Paragraph(f"{readiness} / 100 ({readiness_cat})", body_style)],
        [Paragraph("Expected starting Package", bold_body_style), Paragraph(f"₹{student['salary_low']} - ₹{student['salary_high']} LPA (Average: ₹{student['salary_avg']} LPA)", body_style)]
    ]
    t2 = Table(metrics_data, colWidths=[170, 370])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('PADDING', (0,0), (-1,-1), 4),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
    ]))
    story.append(t2)
    
    reason_style = ParagraphStyle('ReasonStyle', parent=body_style, fontSize=7.5, leading=10, fontName='Helvetica-Oblique')
    story.append(Spacer(1, 3))
    story.append(Paragraph(f"<b>Inference Rationale:</b> {student['prediction_reason']}", reason_style))
    story.append(Spacer(1, 10))
    
    # Section 3: Skill Gap Analysis & Recommendations
    story.append(Paragraph("3. SKILL GAP ANALYSIS & RECOMMENDATIONS", h1_style))
    
    weak_areas = student.get("weak_areas", [])
    if isinstance(weak_areas, str):
        try: weak_areas = json.loads(weak_areas)
        except: weak_areas = []
        
    if weak_areas:
        gap_table_data = [[Paragraph("Area", table_header_style), Paragraph("Score", table_header_style), Paragraph("Action/Recommendation Plan", table_header_style)]]
        for w in weak_areas:
            gap_table_data.append([
                Paragraph(w.get("area", "N/A"), table_cell_style),
                Paragraph(f"{w.get('score', 0)}%", table_cell_style),
                Paragraph(w.get("suggestion", "N/A"), table_cell_style)
            ])
        t3 = Table(gap_table_data, colWidths=[120, 40, 380])
        t3.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6d28d9')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ('PADDING', (0,0), (-1,-1), 4),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8fafc')])
        ]))
        story.append(t3)
    else:
        story.append(Paragraph("No critical skill gaps identified. Student displays exceptional readiness levels.", body_style))
    story.append(Spacer(1, 10))
    
    # Section 4: AI Resume Analysis
    story.append(Paragraph("4. AI RESUME ANALYSIS", h1_style))
    
    res_score = student.get("resume_score", 0)
    missing_skills = student.get("resume_missing_skills", [])
    missing_keywords = student.get("resume_missing_keywords", [])
    strengths_res = student.get("resume_strengths", [])
    suggestions_res = student.get("resume_suggestions", [])
    
    if isinstance(missing_skills, str):
        try: missing_skills = json.loads(missing_skills)
        except: missing_skills = []
    if isinstance(missing_keywords, str):
        try: missing_keywords = json.loads(missing_keywords)
        except: missing_keywords = []
    if isinstance(strengths_res, str):
        try: strengths_res = json.loads(strengths_res)
        except: strengths_res = []
    if isinstance(suggestions_res, str):
        try: suggestions_res = json.loads(suggestions_res)
        except: suggestions_res = []
        
    resume_table_data = [
        [Paragraph("Overall ATS Resume Score", bold_body_style), Paragraph(f"{res_score} / 100", bold_body_style)],
        [Paragraph("Missing Skill Keywords", bold_body_style), Paragraph(", ".join(missing_skills) if missing_skills else "None", body_style)],
        [Paragraph("Missing Core Keywords", bold_body_style), Paragraph(", ".join(missing_keywords) if missing_keywords else "None", body_style)],
        [Paragraph("Key Strengths", bold_body_style), Paragraph(" • " + "\n • ".join(strengths_res) if strengths_res else "None", body_style)],
        [Paragraph("Improvement Suggestions", bold_body_style), Paragraph(" • " + "\n • ".join(suggestions_res) if suggestions_res else "None", body_style)]
    ]
    t4 = Table(resume_table_data, colWidths=[150, 390])
    t4.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
        ('PADDING', (0,0), (-1,-1), 4),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(t4)
    story.append(Spacer(1, 10))
    
    # Section 5: AI Mock Interview Simulator Metrics
    story.append(Paragraph("5. AI MOCK INTERVIEW SIMULATOR METRICS", h1_style))
    
    interview_score = student.get("mock_interview_score", 0)
    int_analysis = student.get("interview_analysis", {})
    int_sug = student.get("interview_suggestions", [])
    
    if isinstance(int_analysis, str):
        try: int_analysis = json.loads(int_analysis)
        except: int_analysis = {}
    if isinstance(int_sug, str):
        try: int_sug = json.loads(int_sug)
        except: int_sug = []
        
    int_strengths = int_analysis.get("strengths", []) if isinstance(int_analysis, dict) else []
    int_weaknesses = int_analysis.get("weaknesses", []) if isinstance(int_analysis, dict) else []
    
    interview_table_data = [
        [Paragraph("Mock Interview Score", bold_body_style), Paragraph(f"{interview_score} / 100", bold_body_style)],
        [Paragraph("Response Strengths", bold_body_style), Paragraph(" • " + "\n • ".join(int_strengths) if int_strengths else "Demonstrates basic department knowledge.", body_style)],
        [Paragraph("Response Weaknesses", bold_body_style), Paragraph(" • " + "\n • ".join(int_weaknesses) if int_weaknesses else "Could improve response completeness and structure.", body_style)],
        [Paragraph("Expert Recommendations", bold_body_style), Paragraph(" • " + "\n • ".join(int_sug) if int_sug else "Practice the STAR method (Situation, Task, Action, Result) for behavioral answers.", body_style)]
    ]
    t5 = Table(interview_table_data, colWidths=[150, 390])
    t5.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
        ('PADDING', (0,0), (-1,-1), 4),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(t5)
    story.append(Spacer(1, 10))
    
    # Section 6: Company Recommendations
    story.append(Paragraph("6. DYNAMIC COMPANY MATCHING RECOMMENDATIONS", h1_style))
    
    rec_companies = student.get("recommended_companies", [])
    if isinstance(rec_companies, str):
        try: rec_companies = json.loads(rec_companies)
        except: rec_companies = []
        
    if rec_companies:
        company_table_data = [[Paragraph("Company", table_header_style), Paragraph("Package Offered", table_header_style), Paragraph("Tier", table_header_style), Paragraph("Match Strength", table_header_style)]]
        for c in rec_companies[:5]:
            company_table_data.append([
                Paragraph(c.get("name", "N/A"), table_cell_style),
                Paragraph(c.get("package", "N/A"), table_cell_style),
                Paragraph(c.get("tier", "N/A"), table_cell_style),
                Paragraph(f"{c.get('match', 0)}% Match", table_cell_style)
            ])
        t6 = Table(company_table_data, colWidths=[140, 140, 120, 140])
        t6.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6d28d9')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ('PADDING', (0,0), (-1,-1), 4),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8fafc')])
        ]))
        story.append(t6)
    else:
        story.append(Paragraph("No matching companies found.", body_style))
    story.append(Spacer(1, 10))
    
    # Section 7: Learning Roadmap
    story.append(Paragraph("7. CURATED 4-MONTH STUDY ROADMAP", h1_style))
    
    roadmap = student.get("learning_roadmap", [])
    if isinstance(roadmap, str):
        try: roadmap = json.loads(roadmap)
        except: roadmap = []
        
    if roadmap:
        road_table_data = [[Paragraph("Timeline", table_header_style), Paragraph("Focus Area", table_header_style), Paragraph("Actionable Milestones", table_header_style)]]
        for r in roadmap[:4]:
            tasks = r.get("tasks", [])
            tasks_str = " • " + "\n • ".join(tasks) if tasks else "N/A"
            road_table_data.append([
                Paragraph(r.get("week", r.get("month", "Month")), table_cell_style),
                Paragraph(r.get("focus", "N/A"), table_cell_style),
                Paragraph(tasks_str, table_cell_style)
            ])
        t7 = Table(road_table_data, colWidths=[80, 120, 340])
        t7.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6d28d9')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ('PADDING', (0,0), (-1,-1), 4),
            ('ALIGN', (0,0), (-1,-1), 'LEFT'),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8fafc')])
        ]))
        story.append(t7)
    else:
        story.append(Paragraph("No custom study roadmap generated.", body_style))
        
    # Build Document
    def add_page_number(canvas, doc):
        canvas.saveState()
        canvas.setFont('Helvetica', 8)
        canvas.setFillColor(colors.HexColor('#64748b'))
        # Draw header rule and text
        canvas.setStrokeColor(colors.HexColor('#cbd5e1'))
        canvas.setLineWidth(0.5)
        canvas.line(36, 756, 576, 756)
        canvas.drawString(36, 762, "AI-Based Placement Prediction and Career Readiness System - Report")
        
        # Draw footer rule and page number
        canvas.line(36, 44, 576, 44)
        page_num = canvas.getPageNumber()
        canvas.drawRightString(576, 32, f"Page {page_num}")
        canvas.drawString(36, 32, f"Report Generated on {time.strftime('%Y-%m-%d %H:%M:%S')} | Confidential")
        canvas.restoreState()

    doc.build(story, onFirstPage=add_page_number, onLaterPages=add_page_number)
    buffer.seek(0)
    return buffer

class InterviewQuestionsRequest(BaseModel):
    register_no: str
    api_key: Optional[str] = None

class InterviewEvaluateRequest(BaseModel):
    register_no: str
    responses: dict # maps question (str) -> response (str)
    api_key: Optional[str] = None

@app.post("/api/interview/questions")
def get_interview_questions(req: InterviewQuestionsRequest):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (req.register_no,))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        raise HTTPException(status_code=404, detail="Student profile not found. Please create prediction profile first.")
        
    student = dict(row)
    name = student["student_name"]
    dept = student["department"]
    cgpa = student["cgpa"]
    prog_skills = student["programming_skills"]
    tech_skills = student["technical_skills"]
    projects = student["projects"]
    internship = student["internship"]
    certifications = student["certifications"]
    prog_langs = student.get("programming_languages", "Python, Java, SQL")
    
    # Check if API Key is supplied or configured
    gemini_key = req.api_key or os.getenv("GEMINI_API_KEY") or "AQ.Ab8RN6JIB7yz38nl6_zsdY3c5S2TzkbHLDui0N0F2C9VUFmHng"
    
    questions = []
    if gemini_key and gemini_key.strip():
        # Call Gemini to dynamically generate questions
        system_instruction = (
            "You are a professional technical recruiter and HR manager. "
            "You output ONLY valid JSON list of strings containing exactly 5 interview questions. "
            "Do NOT include any markdown formatting, backticks, or HTML. Just return the JSON array."
        )
        
        contents = [{
            "role": "user",
            "parts": [{
                "text": (
                    f"Generate exactly 5 interview questions for a student named {name} who is in the {dept} department. "
                    f"Profile details: CGPA is {cgpa}/10, Programming Skill rating is {prog_skills}/100, Technical Score is {tech_skills}/100, "
                    f"has completed {projects} projects, has internship experience: '{internship}', certifications count: {certifications}, "
                    f"and programs in: {prog_langs}.\n\n"
                    "Structure of the 5 questions:\n"
                    "1. HR / Behavioral question\n"
                    "2. General CS/Programming technical question (related to their languages: CS, OOPs, etc.)\n"
                    "3. Core SQL/Database query technical question\n"
                    "4. Department-specific technical question (AIML, ECE, CSE specific)\n"
                    "5. Project / Experience-specific technical question\n\n"
                    "Return ONLY a JSON list of strings, for example: [\"Question 1\", \"Question 2\", \"Question 3\", \"Question 4\", \"Question 5\"]."
                )
            }]
        }]
        
        try:
            reply = call_gemini(gemini_key.strip(), contents, system_instruction)
            # Clean up the JSON if it contains markdown formatting
            cleaned_reply = reply.strip()
            if cleaned_reply.startswith("```json"):
                cleaned_reply = cleaned_reply.split("```json")[1].split("```")[0].strip()
            elif cleaned_reply.startswith("```"):
                cleaned_reply = cleaned_reply.split("```")[1].split("```")[0].strip()
                
            questions = json.loads(cleaned_reply)
            if not isinstance(questions, list) or len(questions) != 5:
                raise ValueError("Parsed result is not a list of 5 elements")
        except Exception as e:
            print(f"Dynamic question generation failed: {e}. Falling back to pre-defined templates.")
            questions = []
            
    # Local fallback questions if Gemini failed or key not present
    if not questions:
        hr_q = f"Hello {name}, tell me about your background, your career aspirations, and why you are interested in software engineering roles."
        
        # Tech 1: Coding / OOP
        tech1_q = f"What is Object-Oriented Programming (OOP)? Explain inheritance, encapsulation, and polymorphism with code examples using {prog_langs}."
        
        # Tech 2: DBMS / SQL
        tech2_q = "Explain the difference between a primary key, unique key, and foreign key. How do INNER JOIN and LEFT OUTER JOIN work in SQL?"
        
        # Tech 3: Department Specific
        if dept == "Artificial Intelligence & Machine Learning":
            dept_q = "Explain the difference between supervised and unsupervised learning. How do you prevent overfitting in classification models (like Decision Trees or Random Forests)?"
        elif dept in ["Computer Science", "Information Technology"]:
            dept_q = "What is the difference between horizontal and vertical scaling in systems design? When should you use a NoSQL database over a relational one?"
        elif dept == "Electronics & Communication":
            dept_q = "What is an RTOS (Real-Time Operating System)? Explain how tasks are scheduled and how semaphores differ from mutexes in embedded firmware."
        elif dept == "Electrical & Electronics":
            dept_q = "What is a PLC (Programmable Logic Controller)? Describe its architecture and how it is programmed for industrial automation using ladder logic."
        else:
            dept_q = "What is Git? Explain how you resolve merge conflicts when collaborating with other developers on a code repository."
            
        # Tech 4: Project
        proj_q = f"Describe one of your {projects} projects. What was the core problem it solved, the technology stack you chose, and the key technical challenge you encountered?"
        
        questions = [hr_q, tech1_q, tech2_q, dept_q, proj_q]
        
    # Save the questions to SQLite for this student
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE students SET 
            interview_questions = ?
        WHERE register_no = ?
    """, (json.dumps(questions), req.register_no))
    conn.commit()
    conn.close()
    
    return {"questions": questions}

@app.post("/api/interview/evaluate")
def evaluate_interview(req: InterviewEvaluateRequest):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (req.register_no,))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        raise HTTPException(status_code=404, detail="Student profile not found")
        
    student = dict(row)
    name = student["student_name"]
    dept = student["department"]
    prog_skills = student["programming_skills"]
    tech_skills = student["technical_skills"]
    
    gemini_key = req.api_key or os.getenv("GEMINI_API_KEY") or "AQ.Ab8RN6JIB7yz38nl6_zsdY3c5S2TzkbHLDui0N0F2C9VUFmHng"
    
    score = 75
    analysis = {
        "strengths": ["Clear definition of concepts", "Structured thought process"],
        "weaknesses": ["Behavioral responses are brief", "Could explain system designs in more detail"]
    }
    suggestions = [
        "Use the STAR method (Situation, Task, Action, Result) for behavioral questions.",
        "Include time/space complexity analysis (Big O) when discussing coding and DSA solutions.",
        "Add details about deployment, database choices, and indexing when discussing projects."
    ]
    
    # Construct prompt content
    qa_pairs = ""
    for idx, (q, r) in enumerate(req.responses.items()):
        qa_pairs += f"Question {idx+1}: {q}\nResponse: {r}\n\n"
        
    if gemini_key and gemini_key.strip():
        system_instruction = (
            "You are an expert AI Technical Recruiter. "
            "You evaluate mock interview answers and output ONLY a valid JSON object. "
            "Do NOT output markdown backticks, explanations, or text. Return ONLY the JSON object."
        )
        
        prompt = (
            f"Evaluate the mock interview responses for the student named {name} (Department: {dept}). "
            "Here are the questions and student's responses:\n\n"
            f"{qa_pairs}\n"
            "Evaluate the quality, correctness, and professional structure of the responses. "
            "Output a JSON object with exactly these fields:\n"
            "{\n"
            "  \"score\": 78, // an integer between 0 and 100\n"
            "  \"analysis\": {\n"
            "    \"strengths\": [\"Strength 1\", \"Strength 2\"], // list of strings, max 3\n"
            "    \"weaknesses\": [\"Weakness 1\", \"Weakness 2\"] // list of strings, max 3\n"
            "  },\n"
            "  \"suggestions\": [\"Suggestion 1\", \"Suggestion 2\"] // list of strings, max 3\n"
            "}\n"
        )
        
        contents = [{
            "role": "user",
            "parts": [{"text": prompt}]
        }]
        
        try:
            reply = call_gemini(gemini_key.strip(), contents, system_instruction)
            cleaned_reply = reply.strip()
            if cleaned_reply.startswith("```json"):
                cleaned_reply = cleaned_reply.split("```json")[1].split("```")[0].strip()
            elif cleaned_reply.startswith("```"):
                cleaned_reply = cleaned_reply.split("```")[1].split("```")[0].strip()
                
            parsed = json.loads(cleaned_reply)
            if "score" in parsed and "analysis" in parsed and "suggestions" in parsed:
                score = int(parsed["score"])
                analysis = parsed["analysis"]
                suggestions = parsed["suggestions"]
        except Exception as e:
            print(f"Dynamic interview evaluation failed: {e}. Falling back to rule-based parser.")
            
    # Local fallback scoring if Gemini fails/key not set
    if gemini_key is None or not gemini_key.strip() or "score" not in locals() or score == 75:
        # Heuristic scoring based on length and core keywords
        word_scores = []
        for q, r in req.responses.items():
            r_clean = r.strip()
            length = len(r_clean)
            
            # Base score for writing something
            q_score = 50
            if length == 0:
                q_score = 10
            elif length < 30:
                q_score = 45 # Very short
            elif length > 120:
                q_score = 80 # Decent length
                
            # Keyword match bonuses
            r_lower = r_clean.lower()
            
            # Check for technical terms
            tech_keywords = ["oop", "class", "object", "polymorphism", "inheritance", "encapsulation", 
                             "join", "table", "select", "key", "index", "sql", "query", "database", 
                             "overfitting", "supervised", "unsupervised", "model", "tree", "learning",
                             "rtos", "mutex", "semaphore", "task", "interrupt", "memory", "complexity", 
                             "git", "conflict", "branch", "commit", "star", "result", "action", "solved"]
                             
            matches = sum(1 for kw in tech_keywords if kw in r_lower)
            q_score += min(matches * 4, 18)
            word_scores.append(min(q_score, 100))
            
        score = int(sum(word_scores) / len(word_scores)) if word_scores else 60
        # Add bounds
        score = max(35, min(score, 94))
        
        # Tailor suggestions based on scores
        if score < 60:
            analysis["weaknesses"].append("Responses are extremely brief and lack technical depth.")
            suggestions.append("Try to expand your responses. Explain the concepts using definitions and code examples where possible.")
        if score >= 80:
            analysis["strengths"].append("Demonstrated solid coding vocabulary and clear conceptual definitions.")
            
    # Save the interview results in database
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE students SET 
            mock_interview_score = ?,
            interview_responses = ?,
            interview_analysis = ?,
            interview_suggestions = ?
        WHERE register_no = ?
    """, (score, json.dumps(req.responses), json.dumps(analysis), json.dumps(suggestions), req.register_no))
    conn.commit()
    
    # IMPORTANT: Fetch updated profile, recalculate readiness score, ML prediction probability and Expected Salary!
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (req.register_no,))
    row = cursor.fetchone()
    
    if row:
        student_data = dict(row)
        # Create StudentProfile object to recalculate
        p = StudentProfile(
            student_name=student_data["student_name"],
            register_no=student_data["register_no"],
            department=student_data["department"],
            gender=student_data["gender"],
            tenth_percentage=float(student_data["tenth_percentage"]),
            twelfth_percentage=float(student_data["twelfth_percentage"]),
            cgpa=float(student_data["cgpa"]),
            backlogs=int(student_data["backlogs"]),
            programming_skills=int(student_data["programming_skills"]),
            aptitude_score=int(student_data["aptitude_score"]),
            communication_skills=int(student_data["communication_skills"]),
            technical_skills=int(student_data["technical_skills"]),
            projects=int(student_data["projects"]),
            internship=student_data["internship"],
            certifications=int(student_data["certifications"]),
            hackathons=int(student_data["hackathons"]),
            resume_uploaded=student_data["resume_uploaded"],
            mock_interview_score=score,
            programming_languages=student_data.get("programming_languages", "Python, Java, SQL")
        )
        
        # Re-run ML predictions with the updated mock interview score
        res = run_predictions(p)
        
        cursor.execute("""
            UPDATE students SET
                placed = ?, probability = ?, readiness_score = ?,
                salary_low = ?, salary_avg = ?, salary_high = ?,
                prediction_reason = ?, weak_areas = ?, learning_roadmap = ?, recommended_companies = ?,
                confidence_score = ?
            WHERE register_no = ?
        """, (
            res["placed"], res["probability"], res["readiness_score"],
            res["salary_low"], res["salary_avg"], res["salary_high"],
            res["prediction_reason"], json.dumps(res["weak_areas"]), json.dumps(res["learning_roadmap"]),
            json.dumps(res["recommended_companies"]), res["confidence_score"], req.register_no
        ))
        conn.commit()
        
    conn.close()
    
    return {
        "score": score,
        "analysis": analysis,
        "suggestions": suggestions
    }

@app.get("/api/pdf-report/{register_no}")
def get_pdf_report(register_no: str):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM students WHERE register_no = ?", (register_no,))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        raise HTTPException(status_code=404, detail="Student profile not found. Please evaluate prediction first.")
        
    student = dict(row)
    pdf_buffer = generate_pdf_report(student)
    
    filename = f"placement_report_{register_no}.pdf"
    
    return StreamingResponse(
        pdf_buffer,
        media_type="application/pdf",
        headers={
            "Content-Disposition": f"attachment; filename={filename}",
            "Cache-Control": "no-cache, no-store, must-revalidate",
            "Pragma": "no-cache"
        }
    )

# ==========================================
# AUTHENTICATION & SESSION MANAGEMENT
# ==========================================
import datetime
import jwt
import bcrypt
from fastapi import Header

JWT_SECRET = "placement_prediction_system_secret_key_2026"
JWT_ALGORITHM = "HS256"

# Pydantic Schemas for Auth
class UserRegisterRequest(BaseModel):
    full_name: str
    username: str
    register_number: str
    email: str
    mobile: str
    password: str
    profile_photo: Optional[str] = None

class UserLoginRequest(BaseModel):
    username_or_email: str
    password: str

class ForgotPasswordRequest(BaseModel):
    email: str

class ResetPasswordRequest(BaseModel):
    email: str
    otp: str
    new_password: str

class ProfileUpdateRequest(BaseModel):
    full_name: Optional[str] = None
    email: Optional[str] = None
    mobile: Optional[str] = None
    password: Optional[str] = None
    profile_photo: Optional[str] = None

# Helper Functions
def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    try:
        return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))
    except Exception:
        return False

def create_jwt_token(user_id: int, username: str, register_number: str) -> str:
    payload = {
        "user_id": user_id,
        "username": username,
        "register_number": register_number,
        "exp": datetime.datetime.utcnow() + datetime.timedelta(days=7)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

def decode_jwt_token(token: str) -> Optional[dict]:
    try:
        return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except Exception:
        return None

# Dependency to fetch logged in user payload
def get_current_user(authorization: Optional[str] = Header(None)):
    if not authorization:
        raise HTTPException(status_code=401, detail="Authorization header missing")
    try:
        parts = authorization.split(" ")
        if len(parts) != 2 or parts[0].lower() != "bearer":
            raise HTTPException(status_code=401, detail="Invalid authorization token format")
        token = parts[1]
        payload = decode_jwt_token(token)
        if not payload:
            raise HTTPException(status_code=401, detail="Session expired or invalid token")
        return payload
    except Exception:
        raise HTTPException(status_code=401, detail="Authorization failed")

# Endpoints
@app.post("/api/register")
def register_user(req: UserRegisterRequest):
    # Validation checks
    # Email validation check
    import re
    if not re.match(r"[^@]+@[^@]+\.[^@]+", req.email):
        raise HTTPException(status_code=400, detail="Invalid email format")
    # Mobile validation check
    if not req.mobile.isdigit() or len(req.mobile) != 10:
        raise HTTPException(status_code=400, detail="Mobile number must be exactly 10 digits")
    # Username check
    if len(req.username) < 4 or len(req.username) > 25 or not re.match(r"^[a-zA-Z0-9_]+$", req.username):
        raise HTTPException(status_code=400, detail="Username must be 4-25 characters (letters, numbers, underscores)")
    # Password complexity check
    # Min 8 chars, 1 uppercase, 1 lowercase, 1 number, 1 special char
    if len(req.password) < 8 or not re.search(r"[A-Z]", req.password) or not re.search(r"[a-z]", req.password) or not re.search(r"[0-9]", req.password) or not re.search(r"[^a-zA-Z0-9]", req.password):
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters, with 1 uppercase, 1 lowercase, 1 digit, and 1 special character")

    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Check duplicate username
    cursor.execute("SELECT id FROM users WHERE username = ?", (req.username,))
    if cursor.fetchone():
        conn.close()
        raise HTTPException(status_code=400, detail="Username is already registered")

    # Check duplicate register_number
    cursor.execute("SELECT id FROM users WHERE register_number = ?", (req.register_number,))
    if cursor.fetchone():
        conn.close()
        raise HTTPException(status_code=400, detail="Register number is already registered")

    # Check duplicate email
    cursor.execute("SELECT id FROM users WHERE email = ?", (req.email,))
    if cursor.fetchone():
        conn.close()
        raise HTTPException(status_code=400, detail="Email is already registered")

    password_hash = hash_password(req.password)
    
    try:
        cursor.execute("""
            INSERT INTO users (full_name, username, register_number, email, mobile, password_hash, profile_photo)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, (req.full_name, req.username, req.register_number, req.email, req.mobile, password_hash, req.profile_photo))
        conn.commit()
    except Exception as e:
        conn.close()
        raise HTTPException(status_code=500, detail=f"Database insertion failed: {e}")
        
    conn.close()
    return {"message": "Registration Successful"}

@app.post("/api/login")
def login_user(req: UserLoginRequest):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute("SELECT * FROM users WHERE username = ? OR email = ?", (req.username_or_email, req.username_or_email))
    user_row = cursor.fetchone()
    
    if not user_row:
        conn.close()
        raise HTTPException(status_code=401, detail="Invalid username or email")
        
    user = dict(user_row)
    if not verify_password(req.password, user["password_hash"]):
        conn.close()
        raise HTTPException(status_code=401, detail="Incorrect password")
        
    # Update last_login
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cursor.execute("UPDATE users SET last_login = ? WHERE id = ?", (now, user["id"]))
    conn.commit()
    conn.close()
    
    token = create_jwt_token(user["id"], user["username"], user["register_number"])
    
    return {
        "token": token,
        "user": {
            "id": user["id"],
            "full_name": user["full_name"],
            "username": user["username"],
            "register_number": user["register_number"],
            "email": user["email"],
            "mobile": user["mobile"],
            "profile_photo": user["profile_photo"],
            "last_login": now
        }
    }

@app.get("/api/profile")
def get_user_profile(current_user: dict = Depends(get_current_user)):
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM users WHERE id = ?", (current_user["user_id"],))
    row = cursor.fetchone()
    conn.close()
    
    if not row:
        raise HTTPException(status_code=404, detail="User not found")
        
    user = dict(row)
    return {
        "id": user["id"],
        "full_name": user["full_name"],
        "username": user["username"],
        "register_number": user["register_number"],
        "email": user["email"],
        "mobile": user["mobile"],
        "profile_photo": user["profile_photo"],
        "created_at": user["created_at"],
        "last_login": user["last_login"]
    }

@app.put("/api/profile")
def update_user_profile(req: ProfileUpdateRequest, current_user: dict = Depends(get_current_user)):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Verify user exists
    cursor.execute("SELECT * FROM users WHERE id = ?", (current_user["user_id"],))
    user_row = cursor.fetchone()
    if not user_row:
        conn.close()
        raise HTTPException(status_code=404, detail="User not found")
    user = dict(user_row)
    
    updates = []
    params = []
    
    if req.full_name is not None:
        updates.append("full_name = ?")
        params.append(req.full_name)
        
    if req.email is not None:
        # Check duplicate email
        cursor.execute("SELECT id FROM users WHERE email = ? AND id != ?", (req.email, current_user["user_id"]))
        if cursor.fetchone():
            conn.close()
            raise HTTPException(status_code=400, detail="Email is already taken by another user")
        updates.append("email = ?")
        params.append(req.email)
        
    if req.mobile is not None:
        if not req.mobile.isdigit() or len(req.mobile) != 10:
            conn.close()
            raise HTTPException(status_code=400, detail="Mobile number must be exactly 10 digits")
        updates.append("mobile = ?")
        params.append(req.mobile)
        
    if req.password is not None and req.password.strip() != "":
        import re
        if len(req.password) < 8 or not re.search(r"[A-Z]", req.password) or not re.search(r"[a-z]", req.password) or not re.search(r"[0-9]", req.password) or not re.search(r"[^a-zA-Z0-9]", req.password):
            conn.close()
            raise HTTPException(status_code=400, detail="Password must be at least 8 characters, with 1 uppercase, 1 lowercase, 1 digit, and 1 special character")
        updates.append("password_hash = ?")
        params.append(hash_password(req.password))
        
    if req.profile_photo is not None:
        updates.append("profile_photo = ?")
        params.append(req.profile_photo)
        
    if updates:
        updates.append("updated_at = CURRENT_TIMESTAMP")
        params.append(current_user["user_id"])
        query = f"UPDATE users SET {', '.join(updates)} WHERE id = ?"
        cursor.execute(query, params)
        conn.commit()
        
        # Also, update student_name in students table if name updated
        if req.full_name is not None:
            cursor.execute("UPDATE students SET student_name = ? WHERE register_no = ?", (req.full_name, current_user["register_number"]))
            conn.commit()
            
    # Fetch updated user
    cursor.execute("SELECT * FROM users WHERE id = ?", (current_user["user_id"],))
    updated_row = cursor.fetchone()
    conn.close()
    
    user = dict(updated_row)
    return {
        "id": user["id"],
        "full_name": user["full_name"],
        "username": user["username"],
        "register_number": user["register_number"],
        "email": user["email"],
        "mobile": user["mobile"],
        "profile_photo": user["profile_photo"]
    }

@app.post("/api/logout")
def logout_user():
    return {"message": "Logout successful"}

@app.post("/api/forgot-password")
def forgot_password(req: ForgotPasswordRequest):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    cursor.execute("SELECT id FROM users WHERE email = ?", (req.email,))
    user = cursor.fetchone()
    if not user:
        conn.close()
        raise HTTPException(status_code=404, detail="Email address is not registered")
        
    # Generate 6-digit OTP
    import random
    otp = str(random.randint(100000, 999999))
    
    # Expiration: 15 mins
    expires_at = (datetime.datetime.now() + datetime.timedelta(minutes=15)).strftime("%Y-%m-%d %H:%M:%S")
    
    cursor.execute("INSERT INTO password_resets (email, otp, expires_at) VALUES (?, ?, ?)", (req.email, otp, expires_at))
    conn.commit()
    conn.close()
    
    # In a real environment we would send an email, but here we return it in response for demo purposes
    return {
        "message": "Verification OTP sent to your registered email.",
        "otp_demo": otp # Return for simulation
    }

@app.post("/api/reset-password")
def reset_password(req: ResetPasswordRequest):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    cursor.execute("SELECT id FROM password_resets WHERE email = ? AND otp = ? AND expires_at > ?", (req.email, req.otp, now))
    reset_row = cursor.fetchone()
    
    if not reset_row:
        conn.close()
        raise HTTPException(status_code=400, detail="Invalid or expired OTP verification code")
        
    # Password complexity check
    import re
    if len(req.new_password) < 8 or not re.search(r"[A-Z]", req.new_password) or not re.search(r"[a-z]", req.new_password) or not re.search(r"[0-9]", req.new_password) or not re.search(r"[^a-zA-Z0-9]", req.new_password):
        conn.close()
        raise HTTPException(status_code=400, detail="Password must be at least 8 characters, with 1 uppercase, 1 lowercase, 1 digit, and 1 special character")
        
    password_hash = hash_password(req.new_password)
    
    cursor.execute("UPDATE users SET password_hash = ? WHERE email = ?", (password_hash, req.email))
    cursor.execute("DELETE FROM password_resets WHERE email = ?", (req.email,))
    conn.commit()
    conn.close()
    
    return {"message": "Password reset successfully"}

@app.get("/api/verify-token")
def verify_token(current_user: dict = Depends(get_current_user)):
    return {"valid": True, "user": current_user}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
